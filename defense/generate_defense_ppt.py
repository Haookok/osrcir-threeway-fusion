"""毕业答辩 PPT 自动生成脚本

生成同济大学本科毕业设计答辩 PPT，约 22 页，16:9 布局。
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------- 配置 ----------------
FIG_DIR = Path("/home/haomingyang03/code/osrcir/docs/thesis/figures")
OUT_DIR = Path("/home/haomingyang03/code/osrcir/defense")
OUT_PPTX = OUT_DIR / "答辩PPT_三路融合.pptx"

# 16:9 尺寸
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 主色调（清爽的深蓝 + 红色辅助）
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x8A)   # 深蓝
COLOR_ACCENT  = RGBColor(0xDC, 0x26, 0x26)   # 红
COLOR_GREEN   = RGBColor(0x16, 0xA3, 0x4A)   # 绿
COLOR_GRAY    = RGBColor(0x4B, 0x55, 0x63)   # 深灰
COLOR_LIGHT   = RGBColor(0xF3, 0xF4, 0xF6)   # 浅灰背景
COLOR_BOX_BG  = RGBColor(0xEE, 0xF2, 0xFF)   # 浅蓝背景

FONT_CN = "微软雅黑"
FONT_EN = "Calibri"


# ---------------- 工具函数 ----------------
def add_text(slide, left, top, width, height, text,
             font_size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name=FONT_CN):
    """添加一个文本框。"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    # 支持多行
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = font_name
        if color is not None:
            run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height,
             fill=None, line=None, line_width=0.5):
    """添加一个圆角矩形作为背景块。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    shape.adjustments[0] = 0.05
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_image(slide, img_path, left, top, width=None, height=None):
    """添加图片。"""
    if width is not None and height is None:
        return slide.shapes.add_picture(str(img_path), left, top, width=width)
    if height is not None and width is None:
        return slide.shapes.add_picture(str(img_path), left, top, height=height)
    return slide.shapes.add_picture(str(img_path), left, top,
                                     width=width, height=height)


def add_title_bar(slide, title_text, subtitle=None):
    """页面顶部的标题栏。"""
    # 顶部蓝色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  0, 0, SLIDE_W, Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    bar.shadow.inherit = False

    # 左侧竖线装饰
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.4), Inches(0.2),
                                     Inches(0.08), Inches(0.45))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_ACCENT
    accent.line.fill.background()
    accent.shadow.inherit = False

    # 标题文字
    add_text(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.55),
             title_text, font_size=24, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF),
             anchor=MSO_ANCHOR.MIDDLE)

    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.3),
                 subtitle, font_size=12,
                 color=RGBColor(0xCC, 0xDB, 0xFE))


def add_page_number(slide, current, total):
    """页脚页码。"""
    add_text(slide, Inches(12.5), Inches(7.1), Inches(0.75), Inches(0.3),
             f"{current} / {total}", font_size=10, color=COLOR_GRAY,
             align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.4), Inches(7.1), Inches(10), Inches(0.3),
             "基于视觉代理与描述融合的零样本组合式图像检索改进  |  杨昊明",
             font_size=10, color=COLOR_GRAY)


# ---------------- 创建 Presentation ----------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# 空白布局
BLANK = prs.slide_layouts[6]


TOTAL_SLIDES = 22


# ==================== 第 1 页: 封面 ====================
slide = prs.slides.add_slide(BLANK)

# 左上装饰块
dec1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0), Inches(0),
                              Inches(4), Inches(0.3))
dec1.fill.solid(); dec1.fill.fore_color.rgb = COLOR_PRIMARY
dec1.line.fill.background(); dec1.shadow.inherit = False

dec2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0), Inches(7.2),
                              Inches(13.333), Inches(0.3))
dec2.fill.solid(); dec2.fill.fore_color.rgb = COLOR_PRIMARY
dec2.line.fill.background(); dec2.shadow.inherit = False

dec3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0), Inches(7.15),
                              Inches(13.333), Inches(0.05))
dec3.fill.solid(); dec3.fill.fore_color.rgb = COLOR_ACCENT
dec3.line.fill.background(); dec3.shadow.inherit = False

add_text(slide, Inches(0.5), Inches(0.7), Inches(12.333), Inches(0.5),
         "TONGJI UNIVERSITY", font_size=20, bold=True,
         color=COLOR_PRIMARY, align=PP_ALIGN.CENTER, font_name=FONT_EN)
add_text(slide, Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.5),
         "同济大学本科毕业设计（论文）答辩",
         font_size=20, bold=True, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

# 主标题
add_text(slide, Inches(0.5), Inches(2.6), Inches(12.333), Inches(1.0),
         "基于视觉代理与描述融合的",
         font_size=40, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.5), Inches(12.333), Inches(1.0),
         "零样本组合式图像检索改进",
         font_size=40, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)

# 副标题
add_text(slide, Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.5),
         "— 三路融合：视觉代理 + 反幻觉精炼 + 描述集成 —",
         font_size=20, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)

# 下部信息
add_text(slide, Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.4),
         "学生姓名：杨昊明    |    学院：电子与信息工程学院",
         font_size=16, color=COLOR_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(6.25), Inches(12.333), Inches(0.4),
         "专业：计算机科学与技术    |    指导教师：（待填写）",
         font_size=16, color=COLOR_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4),
         "2026 年 5 月",
         font_size=14, color=COLOR_GRAY, align=PP_ALIGN.CENTER)


# ==================== 第 2 页: 目录 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "目录", "Contents")

contents = [
    ("01", "研究背景与问题", "组合式图像检索与零样本挑战"),
    ("02", "相关工作与基线分析", "OSrCIR 方法回顾与不足"),
    ("03", "三路融合方法", "核心思路与公式推导"),
    ("04", "实验设计与主要结果", "9 个数据集全量验证"),
    ("05", "消融实验与参数分析", "模块贡献与 α/β 网格搜索"),
    ("06", "总结与展望", "主要贡献与未来工作"),
]

for i, (num, zh, en) in enumerate(contents):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(1.5 + row * 1.7)

    # 编号
    add_text(slide, left, top, Inches(1.2), Inches(1.2),
             num, font_size=54, bold=True, color=COLOR_ACCENT,
             anchor=MSO_ANCHOR.MIDDLE, font_name=FONT_EN)
    # 中文标题
    add_text(slide, left + Inches(1.3), top + Inches(0.1),
             Inches(5), Inches(0.5),
             zh, font_size=22, bold=True, color=COLOR_PRIMARY)
    # 说明
    add_text(slide, left + Inches(1.3), top + Inches(0.6),
             Inches(5), Inches(0.4),
             en, font_size=14, color=COLOR_GRAY)

add_page_number(slide, 2, TOTAL_SLIDES)


# ==================== 第 3 页: 研究背景 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "01  研究背景",
              "Composed Image Retrieval:  图像 + 文本 → 目标图像")

# 左侧文字
add_text(slide, Inches(0.6), Inches(1.2), Inches(6), Inches(0.4),
         "▎ 任务定义", font_size=18, bold=True, color=COLOR_PRIMARY)
add_text(slide, Inches(0.6), Inches(1.7), Inches(6), Inches(1.8),
         "组合式图像检索 (Composed Image Retrieval, CIR)：\n"
         "给定一张参考图像和一段修改文本，\n"
         "从图库中检索符合修改意图的目标图像。",
         font_size=14, color=COLOR_GRAY)

add_text(slide, Inches(0.6), Inches(3.3), Inches(6), Inches(0.4),
         "▎ 零样本设定 (Zero-Shot)", font_size=18, bold=True, color=COLOR_PRIMARY)
add_text(slide, Inches(0.6), Inches(3.8), Inches(6), Inches(1.8),
         "不依赖目标数据集的监督训练\n"
         "仅使用预训练视觉语言模型完成检索\n"
         "更贴近实际部署中的跨域应用需求",
         font_size=14, color=COLOR_GRAY)

add_text(slide, Inches(0.6), Inches(5.4), Inches(6), Inches(0.4),
         "▎ 应用价值", font_size=18, bold=True, color=COLOR_PRIMARY)
add_text(slide, Inches(0.6), Inches(5.9), Inches(6), Inches(1.0),
         "电商同款搜索 · 图像编辑辅助 · 内容创作工具\n"
         "跨模态信息检索 · 智能推荐系统",
         font_size=14, color=COLOR_GRAY)

# 右侧示例框
add_rect(slide, Inches(7.2), Inches(1.3), Inches(5.6), Inches(5.4),
         fill=COLOR_BOX_BG, line=COLOR_PRIMARY, line_width=1)

add_text(slide, Inches(7.4), Inches(1.5), Inches(5.2), Inches(0.4),
         "▎ 典型示例", font_size=16, bold=True, color=COLOR_PRIMARY)

add_text(slide, Inches(7.4), Inches(2.1), Inches(5.2), Inches(0.5),
         "服饰检索", font_size=14, bold=True, color=COLOR_ACCENT)
add_text(slide, Inches(7.4), Inches(2.55), Inches(5.2), Inches(1),
         "输入：一张黑色连衣裙图片\n"
         "文本：\"change the color to red\"\n"
         "输出：款式相同但颜色为红色的连衣裙",
         font_size=13, color=COLOR_GRAY)

add_text(slide, Inches(7.4), Inches(4.0), Inches(5.2), Inches(0.5),
         "自然场景", font_size=14, bold=True, color=COLOR_ACCENT)
add_text(slide, Inches(7.4), Inches(4.45), Inches(5.2), Inches(1.3),
         "输入：街道上有一辆汽车\n"
         "文本：\"replace the car with a bus\"\n"
         "输出：同样街景但车辆变为公交车",
         font_size=13, color=COLOR_GRAY)

add_text(slide, Inches(7.4), Inches(5.9), Inches(5.2), Inches(0.4),
         "⚠ 难点：既要保留参考图主体，又要准确应用文本修改",
         font_size=12, bold=True, color=COLOR_ACCENT)

add_page_number(slide, 3, TOTAL_SLIDES)


# ==================== 第 4 页: 研究现状与基线 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "01  研究现状", "零样本组合式图像检索的两条技术路线")

# 左侧：基于映射的方法
add_rect(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5),
         fill=COLOR_LIGHT, line=COLOR_GRAY, line_width=0.5)
add_text(slide, Inches(0.75), Inches(1.4), Inches(5.5), Inches(0.4),
         "▎ 基于映射的方法", font_size=18, bold=True, color=COLOR_PRIMARY)
add_text(slide, Inches(0.75), Inches(1.9), Inches(5.5), Inches(0.3),
         "Mapping-Based Approaches", font_size=11,
         color=COLOR_GRAY, font_name=FONT_EN)

add_text(slide, Inches(0.75), Inches(2.4), Inches(5.5), Inches(0.4),
         "核心思路：", font_size=13, bold=True, color=COLOR_GRAY)
add_text(slide, Inches(0.75), Inches(2.8), Inches(5.5), Inches(0.5),
         "把图像映射到文本空间的伪词向量，与修改文本组合后编码",
         font_size=12, color=COLOR_GRAY)

add_text(slide, Inches(0.75), Inches(3.6), Inches(5.5), Inches(0.4),
         "代表工作：", font_size=13, bold=True, color=COLOR_GRAY)
add_text(slide, Inches(0.75), Inches(4.0), Inches(5.5), Inches(1.0),
         "• Pic2Word（CVPR 2023）\n"
         "• SEARLE / iSEARLE（ICCV 2023）\n"
         "• LinCIR（CVPR 2024）",
         font_size=12, color=COLOR_GRAY)

add_text(slide, Inches(0.75), Inches(5.3), Inches(5.5), Inches(0.4),
         "局限：", font_size=13, bold=True, color=COLOR_ACCENT)
add_text(slide, Inches(0.75), Inches(5.7), Inches(5.5), Inches(1.0),
         "图像特征被压缩为文本向量时丢失\n细粒度视觉信息",
         font_size=12, color=COLOR_GRAY)

# 右侧：基于推理的方法
add_rect(slide, Inches(6.9), Inches(1.2), Inches(6), Inches(5.5),
         fill=COLOR_BOX_BG, line=COLOR_PRIMARY, line_width=1)
add_text(slide, Inches(7.15), Inches(1.4), Inches(5.5), Inches(0.4),
         "▎ 基于推理的方法（本文基线）", font_size=18,
         bold=True, color=COLOR_PRIMARY)
add_text(slide, Inches(7.15), Inches(1.9), Inches(5.5), Inches(0.3),
         "Reasoning-Based Approaches", font_size=11,
         color=COLOR_GRAY, font_name=FONT_EN)

add_text(slide, Inches(7.15), Inches(2.4), Inches(5.5), Inches(0.4),
         "核心思路：", font_size=13, bold=True, color=COLOR_GRAY)
add_text(slide, Inches(7.15), Inches(2.8), Inches(5.5), Inches(0.5),
         "用 MLLM 直接推理目标图像的文本描述，再经 CLIP 检索",
         font_size=12, color=COLOR_GRAY)

add_text(slide, Inches(7.15), Inches(3.6), Inches(5.5), Inches(0.4),
         "代表工作：", font_size=13, bold=True, color=COLOR_GRAY)
add_text(slide, Inches(7.15), Inches(4.0), Inches(5.5), Inches(1.0),
         "• CIReVL（ICLR 2024）：两阶段描述推理\n"
         "• OSrCIR（CVPR 2025 Highlight）：单阶段\n"
         "   反思式思维链 — 本文基线",
         font_size=12, color=COLOR_GRAY)

add_text(slide, Inches(7.15), Inches(5.3), Inches(5.5), Inches(0.4),
         "遗留问题：", font_size=13, bold=True, color=COLOR_ACCENT)
add_text(slide, Inches(7.15), Inches(5.7), Inches(5.5), Inches(1.0),
         "纯文本推理路径，缺乏视觉验证机制\n单点失误难以自我纠正",
         font_size=12, color=COLOR_GRAY)

add_page_number(slide, 4, TOTAL_SLIDES)


# ==================== 第 5 页: 问题与思路 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "01  现有方法的三个问题  →  本文的三个创新",
              "Problems & Motivations")

problems = [
    ("① 初始描述缺乏视觉验证",
     "MLLM 只调用一次，描述一旦偏差便无法纠正",
     "创新1：视觉代理 (Visual Proxy)",
     "把描述还原为图像，提供视觉校验与图像空间信号"),
    ("② 单一文本路径过于脆弱",
     "CLIP 对文本表述敏感，单描述易受随机波动影响",
     "创新3：描述融合 (Ensemble)",
     "CLIP 特征空间加权融合 D₁ 与 D₂，提升稳健性"),
    ("③ AI 图像引入幻觉风险",
     "直接描述代理图会把虚构细节带入目标描述",
     "创新2：反幻觉提示词 (V7 Prompt)",
     "代理图仅作诊断工具，禁止复制其视觉细节"),
]

for i, (p_title, p_desc, s_title, s_desc) in enumerate(problems):
    top = Inches(1.15 + i * 1.9)

    # 问题侧
    add_rect(slide, Inches(0.5), top, Inches(5.8), Inches(1.7),
             fill=RGBColor(0xFE, 0xE2, 0xE2), line=COLOR_ACCENT, line_width=1)
    add_text(slide, Inches(0.75), top + Inches(0.15), Inches(5.4), Inches(0.5),
             p_title, font_size=16, bold=True, color=COLOR_ACCENT)
    add_text(slide, Inches(0.75), top + Inches(0.75), Inches(5.4), Inches(0.9),
             p_desc, font_size=13, color=COLOR_GRAY)

    # 箭头
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(6.4), top + Inches(0.55),
                                    Inches(0.6), Inches(0.5))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = COLOR_PRIMARY
    arrow.line.fill.background(); arrow.shadow.inherit = False

    # 方案侧
    add_rect(slide, Inches(7.1), top, Inches(5.8), Inches(1.7),
             fill=RGBColor(0xDC, 0xFC, 0xE7), line=COLOR_GREEN, line_width=1)
    add_text(slide, Inches(7.35), top + Inches(0.15), Inches(5.4), Inches(0.5),
             s_title, font_size=16, bold=True, color=COLOR_GREEN)
    add_text(slide, Inches(7.35), top + Inches(0.75), Inches(5.4), Inches(0.9),
             s_desc, font_size=13, color=COLOR_GRAY)

add_page_number(slide, 5, TOTAL_SLIDES)


# ==================== 第 6 页: 三路融合总体流程 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "02  三路融合方法 — 总体流程",
              "Three-Way Fusion Framework")

# 插入流程图
add_image(slide, FIG_DIR / "fig_pipeline.png",
          Inches(0.7), Inches(1.1), width=Inches(12))

# 底部要点
add_text(slide, Inches(0.7), Inches(6.25), Inches(12), Inches(0.4),
         "▎ 核心思想：文本推理 + 视觉代理 + 精炼描述 三信号互补，告别单一路径",
         font_size=14, bold=True, color=COLOR_PRIMARY)
add_text(slide, Inches(0.7), Inches(6.65), Inches(12), Inches(0.35),
         "推理成本增加：多 1 次文生图 + 1 次 MLLM 精炼；但全程无需训练，保留零样本通用性",
         font_size=12, color=COLOR_GRAY)

add_page_number(slide, 6, TOTAL_SLIDES)


# ==================== 第 7 页: 创新1 — 视觉代理 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "02  创新 1 — 视觉代理 (Visual Proxy)",
              "用文生图模型把描述 D₁ 还原为图像，引入视觉空间信号")

# 左侧：流程说明
add_rect(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5),
         fill=COLOR_BOX_BG, line=COLOR_PRIMARY, line_width=0.8)

add_text(slide, Inches(0.75), Inches(1.4), Inches(5.5), Inches(0.4),
         "▎ 执行流程", font_size=18, bold=True, color=COLOR_PRIMARY)

add_text(slide, Inches(0.75), Inches(2.0), Inches(5.5), Inches(0.5),
         "Step 1", font_size=12, bold=True, color=COLOR_ACCENT)
add_text(slide, Inches(0.75), Inches(2.4), Inches(5.5), Inches(0.7),
         "MLLM (Qwen-VL-Max) 生成初始描述 D₁",
         font_size=13, color=COLOR_GRAY)

add_text(slide, Inches(0.75), Inches(3.1), Inches(5.5), Inches(0.5),
         "Step 2", font_size=12, bold=True, color=COLOR_ACCENT)
add_text(slide, Inches(0.75), Inches(3.5), Inches(5.5), Inches(0.7),
         "文生图模型 (MiniMax image-01) 将 D₁ 渲染为代理图像 P",
         font_size=13, color=COLOR_GRAY)

add_text(slide, Inches(0.75), Inches(4.3), Inches(5.5), Inches(0.5),
         "Step 3", font_size=12, bold=True, color=COLOR_ACCENT)
add_text(slide, Inches(0.75), Inches(4.7), Inches(5.5), Inches(1.5),
         "代理图 P 提供两类作用：\n"
         "  · 检索阶段贡献图像空间相似度信号\n"
         "  · 为第二轮精炼提供视觉对照依据",
         font_size=13, color=COLOR_GRAY)

# 右侧：作用分析
add_rect(slide, Inches(6.9), Inches(1.2), Inches(6), Inches(5.5),
         fill=COLOR_LIGHT, line=COLOR_GRAY, line_width=0.5)

add_text(slide, Inches(7.15), Inches(1.4), Inches(5.5), Inches(0.4),
         "▎ 小规模验证 (FashionIQ dress, 50 样本)",
         font_size=16, bold=True, color=COLOR_PRIMARY)

# 迷你表格
header_top = Inches(2.1)
cols = [("方法", 3.2), ("R@10", 1.8)]
left_base = Inches(7.3)
row_h = Inches(0.45)

# 表头
header = [("Baseline（纯 D₁）", "18.0"),
          ("Plan A（代理图后融合，α=0.8）", "26.0"),
          ("Plan B（代理图前融合送 MLLM）", "22.0"),
          ("🏆 三路融合（A+B 组合）", "最终方案")]

add_rect(slide, left_base, header_top, Inches(5.3), row_h,
         fill=COLOR_PRIMARY, line=None)
add_text(slide, left_base + Inches(0.1), header_top + Inches(0.08),
         Inches(3.2), Inches(0.35), "方法", font_size=13,
         bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(slide, left_base + Inches(3.5), header_top + Inches(0.08),
         Inches(1.7), Inches(0.35), "dress R@10", font_size=13,
         bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

for i, (m, v) in enumerate(header):
    rt = header_top + row_h * (i + 1)
    bg = COLOR_LIGHT if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
    add_rect(slide, left_base, rt, Inches(5.3), row_h,
             fill=bg, line=RGBColor(0xE5, 0xE7, 0xEB), line_width=0.25)
    add_text(slide, left_base + Inches(0.1), rt + Inches(0.08),
             Inches(3.4), Inches(0.35), m, font_size=12, color=COLOR_GRAY)
    color = COLOR_ACCENT if i >= 1 else COLOR_GRAY
    bold = i >= 1
    add_text(slide, left_base + Inches(3.5), rt + Inches(0.08),
             Inches(1.7), Inches(0.35), v, font_size=12,
             bold=bold, color=color, align=PP_ALIGN.CENTER)

add_text(slide, Inches(7.15), Inches(5.2), Inches(5.5), Inches(0.5),
         "▎ 结论", font_size=16, bold=True, color=COLOR_PRIMARY)
add_text(slide, Inches(7.15), Inches(5.7), Inches(5.5), Inches(1.0),
         "代理图虽含噪声，但 +8 pp 的显著提升\n证明其检索价值真实存在",
         font_size=13, color=COLOR_GRAY)

add_page_number(slide, 7, TOTAL_SLIDES)


# ==================== 第 8 页: 创新2 — 反幻觉精炼 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "02  创新 2 — 反幻觉精炼 (V7 Anti-Hallucination CoT)",
              "代理图 \"只做诊断工具，不做描述来源\"")

# 上方问题
add_rect(slide, Inches(0.5), Inches(1.15), Inches(12.333), Inches(1.2),
         fill=RGBColor(0xFE, 0xE2, 0xE2), line=COLOR_ACCENT, line_width=1)
add_text(slide, Inches(0.75), Inches(1.3), Inches(12), Inches(0.4),
         "❗ 问题：直接让 MLLM 对比代理图会引入严重幻觉",
         font_size=16, bold=True, color=COLOR_ACCENT)
add_text(slide, Inches(0.75), Inches(1.75), Inches(12), Inches(0.6),
         "代理图是 AI 生成产物，其背景/材质/环境细节均为虚构。\n"
         "若模型把这些细节写入目标描述，检索反而偏离真实目标（早期实验 CIRR R@10 -17.5）。",
         font_size=13, color=COLOR_GRAY)

# 下方：V7 策略三原则
add_text(slide, Inches(0.5), Inches(2.6), Inches(12), Inches(0.4),
         "▎ V7 反幻觉策略的三条核心原则", font_size=18,
         bold=True, color=COLOR_PRIMARY)

rules = [
    ("01", "明确声明代理图是 AI 生成",
     "在 prompt 中告诉 MLLM 代理图不可信，可能含幻觉细节"),
    ("02", "代理图仅用于诊断性检查",
     "只检查 D₁ 是否正确应用了修改（颜色/物体/属性）"),
    ("03", "强制输出简短描述",
     "只保留核心属性，不超过修改文本长度（后针对 GeneCIS 放宽）"),
]

for i, (num, title, desc) in enumerate(rules):
    left = Inches(0.5 + i * 4.3)
    add_rect(slide, left, Inches(3.2), Inches(4.1), Inches(2.6),
             fill=COLOR_BOX_BG, line=COLOR_PRIMARY, line_width=1)
    add_text(slide, left + Inches(0.25), Inches(3.35),
             Inches(1.0), Inches(0.8),
             num, font_size=36, bold=True, color=COLOR_ACCENT,
             font_name=FONT_EN)
    add_text(slide, left + Inches(1.35), Inches(3.45),
             Inches(2.6), Inches(0.8),
             title, font_size=14, bold=True, color=COLOR_PRIMARY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, left + Inches(0.25), Inches(4.4),
             Inches(3.7), Inches(1.3),
             desc, font_size=12, color=COLOR_GRAY)

# 底部核心变化
add_rect(slide, Inches(0.5), Inches(6.05), Inches(12.333), Inches(0.85),
         fill=RGBColor(0xDC, 0xFC, 0xE7), line=COLOR_GREEN, line_width=1)
add_text(slide, Inches(0.75), Inches(6.15), Inches(12), Inches(0.4),
         "🎯 任务重新定义：", font_size=14, bold=True, color=COLOR_GREEN)
add_text(slide, Inches(0.75), Inches(6.5), Inches(12), Inches(0.4),
         "从「重新描写目标图像」 →  变为  「在 D₁ 基础上做最小必要修正」",
         font_size=13, color=COLOR_GRAY)

add_page_number(slide, 8, TOTAL_SLIDES)


# ==================== 第 9 页: Prompt 演进过程 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "02  Prompt 演进过程 — 从失败中突破",
              "6 个版本迭代，关键洞察：限制信息来源比增加信息更重要")

# 插入 prompt 演进图
add_image(slide, FIG_DIR / "fig_prompt_evolution.png",
          Inches(0.5), Inches(1.1), height=Inches(4.6))

# 右侧要点
add_text(slide, Inches(7.5), Inches(1.3), Inches(5.5), Inches(0.4),
         "▎ 关键发现", font_size=18, bold=True, color=COLOR_PRIMARY)

insights = [
    ("V6 灾难性退化",
     "把 D₁ 完整上下文送入 MLLM → R@10 跌到 49.5 (-17.5)",
     COLOR_ACCENT),
    ("V7 突破",
     "切断「从代理图复制细节」的路径 → +6.5",
     COLOR_GREEN),
    ("V7 + Ensemble",
     "牺牲极端提升换取全数据集稳定（见创新3）",
     COLOR_PRIMARY),
]

for i, (t, d, c) in enumerate(insights):
    top = Inches(2.0 + i * 1.3)
    add_text(slide, Inches(7.5), top, Inches(5.5), Inches(0.4),
             f"• {t}", font_size=14, bold=True, color=c)
    add_text(slide, Inches(7.5), top + Inches(0.45), Inches(5.5), Inches(0.8),
             d, font_size=12, color=COLOR_GRAY)

# 底部总结
add_rect(slide, Inches(0.5), Inches(6.05), Inches(12.333), Inches(0.85),
         fill=COLOR_LIGHT, line=COLOR_GRAY, line_width=0.5)
add_text(slide, Inches(0.75), Inches(6.15), Inches(12), Inches(0.4),
         "💡 核心教训：",
         font_size=13, bold=True, color=COLOR_PRIMARY)
add_text(slide, Inches(0.75), Inches(6.5), Inches(12), Inches(0.4),
         "MLLM 在多模态场景下的幻觉风险要靠 prompt 结构约束，而非依赖模型自身能力",
         font_size=12, color=COLOR_GRAY)

add_page_number(slide, 9, TOTAL_SLIDES)


# ==================== 第 10 页: 创新3 — 描述融合 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "02  创新 3 — 描述融合 (Description Ensemble)",
              "在 CLIP 特征空间加权融合 D₁ 与 D₂，防止退化")

# 左侧：动机
add_rect(slide, Inches(0.5), Inches(1.15), Inches(6.2), Inches(5.7),
         fill=COLOR_LIGHT, line=COLOR_GRAY, line_width=0.5)

add_text(slide, Inches(0.75), Inches(1.3), Inches(5.8), Inches(0.4),
         "▎ 为什么不直接用 D₂ 替换 D₁？",
         font_size=16, bold=True, color=COLOR_PRIMARY)

add_text(slide, Inches(0.75), Inches(1.9), Inches(5.8), Inches(1.6),
         "V7 精炼虽然对多数样本有益，但：\n"
         "  · 部分样本存在信息过度压缩\n"
         "  · FashionIQ shirt 子集会退化 (-4.0)\n"
         "  · 「全替换」是「全有或全无」决策",
         font_size=13, color=COLOR_GRAY)

add_text(slide, Inches(0.75), Inches(3.6), Inches(5.8), Inches(0.4),
         "▎ 融合思想",
         font_size=16, bold=True, color=COLOR_PRIMARY)
add_text(slide, Inches(0.75), Inches(4.0), Inches(5.8), Inches(2.5),
         "不替换，改为在 CLIP 特征空间加权平均：\n"
         "即使 D₂ 在个别样本上失败，\n"
         "D₁ 的信号仍占主导地位 → 整体稳健",
         font_size=13, color=COLOR_GRAY)

# 右侧：公式
add_rect(slide, Inches(7.0), Inches(1.15), Inches(5.9), Inches(5.7),
         fill=COLOR_BOX_BG, line=COLOR_PRIMARY, line_width=1)

add_text(slide, Inches(7.25), Inches(1.3), Inches(5.5), Inches(0.4),
         "▎ 融合公式",
         font_size=16, bold=True, color=COLOR_PRIMARY)

add_text(slide, Inches(7.25), Inches(2.0), Inches(5.5), Inches(1.0),
         "f_text = normalize(",
         font_size=16, color=COLOR_GRAY, font_name=FONT_EN)
add_text(slide, Inches(7.25), Inches(2.4), Inches(5.5), Inches(0.6),
         "   β · CLIP(D₁) + (1-β) · CLIP(D₂)  )",
         font_size=16, color=COLOR_ACCENT,
         bold=True, font_name=FONT_EN)

# 关键参数
add_text(slide, Inches(7.25), Inches(3.5), Inches(5.5), Inches(0.5),
         "🔑 默认 β = 0.7", font_size=20,
         bold=True, color=COLOR_GREEN)
add_text(slide, Inches(7.25), Inches(4.1), Inches(5.5), Inches(1.8),
         "  含义：\n"
         "  · D₁（原始描述）占 70%\n"
         "  · D₂（精炼描述）占 30%\n"
         "  · 以 D₁ 为主，D₂ 辅助修正",
         font_size=13, color=COLOR_GRAY)

add_text(slide, Inches(7.25), Inches(5.9), Inches(5.5), Inches(0.8),
         "⚙ β 是可调超参：\n"
         "  GeneCIS 上可根据任务切换 (0.30 ~ 1.00)",
         font_size=12, color=COLOR_GRAY)

add_page_number(slide, 10, TOTAL_SLIDES)


# ==================== 第 11 页: 三路融合公式 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "02  三路融合最终公式",
              "文本融合 + 代理图 = 三路信号互补")

# 公式区
add_rect(slide, Inches(0.5), Inches(1.2), Inches(12.333), Inches(2.8),
         fill=COLOR_BOX_BG, line=COLOR_PRIMARY, line_width=1)

add_text(slide, Inches(0.75), Inches(1.4), Inches(12), Inches(0.4),
         "▎ 最终检索得分",
         font_size=18, bold=True, color=COLOR_PRIMARY)

add_text(slide, Inches(0.75), Inches(2.1), Inches(12), Inches(0.7),
         "score(I_c) =",
         font_size=22, bold=True, color=COLOR_GRAY, font_name=FONT_EN)
add_text(slide, Inches(2.8), Inches(2.1), Inches(10), Inches(0.7),
         "α · sim(f_text, CLIP(I_c))  +  (1-α) · sim(CLIP(P), CLIP(I_c))",
         font_size=22, bold=True, color=COLOR_ACCENT, font_name=FONT_EN)

add_text(slide, Inches(0.75), Inches(3.0), Inches(12), Inches(0.4),
         "其中 f_text = normalize(β · CLIP(D₁) + (1-β) · CLIP(D₂))",
         font_size=14, color=COLOR_GRAY, font_name=FONT_EN)

add_text(slide, Inches(0.75), Inches(3.45), Inches(12), Inches(0.4),
         "默认参数：α = 0.9（文本为主）, β = 0.7（原始描述为主）",
         font_size=13, color=COLOR_GRAY)

# 参数说明三列
add_text(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(0.4),
         "▎ 超参数的物理含义", font_size=16, bold=True, color=COLOR_PRIMARY)

param_boxes = [
    ("β", "控制 D₁ 与 D₂ 的权重",
     "β=1.0  只用原始描述（= OSrCIR）\n"
     "β=0.0  只用精炼描述（激进）\n"
     "β=0.7  默认：以 D₁ 兜底，D₂ 修正", COLOR_GREEN),
    ("α", "控制文本 vs 代理图权重",
     "α=1.0  只用文本检索\n"
     "α=0.0  只用代理图检索（过度冒险）\n"
     "α=0.9  默认：文本为主，图像辅助", COLOR_PRIMARY),
    ("组合", "默认 (β=0.7, α=0.9)",
     "FashionIQ/CIRCO/CIRR 使用默认\n"
     "GeneCIS 用 task-adaptive 调优\n"
     "后续可学习自适应预测", COLOR_ACCENT),
]

for i, (sym, desc, detail, c) in enumerate(param_boxes):
    left = Inches(0.5 + i * 4.3)
    add_rect(slide, left, Inches(4.85), Inches(4.1), Inches(2.1),
             fill=COLOR_LIGHT, line=c, line_width=1)
    add_text(slide, left + Inches(0.2), Inches(4.95),
             Inches(1.0), Inches(0.6),
             sym, font_size=32, bold=True, color=c, font_name=FONT_EN)
    add_text(slide, left + Inches(1.3), Inches(5.05),
             Inches(2.7), Inches(0.6),
             desc, font_size=12, bold=True, color=COLOR_GRAY)
    add_text(slide, left + Inches(0.2), Inches(5.7),
             Inches(3.8), Inches(1.2),
             detail, font_size=11, color=COLOR_GRAY)

add_page_number(slide, 11, TOTAL_SLIDES)


# ==================== 第 12 页: 实验设置 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "03  实验设置", "9 个数据集 · 统一基线 · 公平对比")

# 左上：数据集
add_rect(slide, Inches(0.5), Inches(1.15), Inches(6.2), Inches(3.1),
         fill=COLOR_BOX_BG, line=COLOR_PRIMARY, line_width=0.8)
add_text(slide, Inches(0.75), Inches(1.3), Inches(5.8), Inches(0.4),
         "▎ 数据集（9 个子任务）", font_size=16, bold=True, color=COLOR_PRIMARY)

ds_info = [
    ("FashionIQ", "dress / shirt / toptee（服饰）", "3 子集"),
    ("CIRCO", "开放域 (gallery ≈ 123K)", "1 子集"),
    ("CIRR", "自然场景 (gallery ≈ 2.3K)", "1 子集"),
    ("GeneCIS", "change/focus × object/attribute", "4 子集"),
]
for i, (name, desc, num) in enumerate(ds_info):
    top = Inches(1.85 + i * 0.55)
    add_text(slide, Inches(0.85), top, Inches(1.8), Inches(0.4),
             f"• {name}", font_size=13, bold=True, color=COLOR_PRIMARY)
    add_text(slide, Inches(2.7), top, Inches(3.3), Inches(0.4),
             desc, font_size=12, color=COLOR_GRAY)
    add_text(slide, Inches(5.8), top, Inches(0.9), Inches(0.4),
             num, font_size=12, color=COLOR_ACCENT, align=PP_ALIGN.RIGHT)

# 右上：实现细节
add_rect(slide, Inches(6.9), Inches(1.15), Inches(6), Inches(3.1),
         fill=COLOR_LIGHT, line=COLOR_GRAY, line_width=0.5)
add_text(slide, Inches(7.15), Inches(1.3), Inches(5.6), Inches(0.4),
         "▎ 实现细节", font_size=16, bold=True, color=COLOR_PRIMARY)

detail_info = [
    ("MLLM", "Qwen-VL-Max (阿里云 DashScope)"),
    ("文生图模型", "MiniMax image-01"),
    ("视觉编码器", "CLIP ViT-L/14 (OpenAI)"),
    ("默认参数", "β=0.7, α=0.9"),
    ("评估指标", "Recall@K / mAP@K"),
]
for i, (k, v) in enumerate(detail_info):
    top = Inches(1.85 + i * 0.45)
    add_text(slide, Inches(7.25), top, Inches(1.8), Inches(0.4),
             f"· {k}", font_size=13, bold=True, color=COLOR_PRIMARY)
    add_text(slide, Inches(9.05), top, Inches(3.6), Inches(0.4),
             v, font_size=12, color=COLOR_GRAY)

# 下部：基线说明
add_rect(slide, Inches(0.5), Inches(4.45), Inches(12.333), Inches(2.4),
         fill=RGBColor(0xFE, 0xF3, 0xC7), line=RGBColor(0xD9, 0x77, 0x06),
         line_width=1)
add_text(slide, Inches(0.75), Inches(4.6), Inches(12), Inches(0.4),
         "⚠ 基线复现说明（关键公平性问题）",
         font_size=16, bold=True, color=RGBColor(0xB4, 0x53, 0x09))

add_text(slide, Inches(0.75), Inches(5.1), Inches(12), Inches(0.5),
         "原论文 OSrCIR 使用 GPT-4o；本工作因接口可用性限制改用 Qwen-VL-Max。",
         font_size=13, color=COLOR_GRAY)
add_text(slide, Inches(0.75), Inches(5.55), Inches(12), Inches(0.5),
         "因此 Baseline 数值低于原论文（如 FIQ dress R@10: 29.70 → 15.80），",
         font_size=13, color=COLOR_GRAY)
add_text(slide, Inches(0.75), Inches(6.0), Inches(12), Inches(0.5),
         "但 所有改进实验均基于相同 Baseline 进行对比，改进幅度的比较仍然公平有效。",
         font_size=13, bold=True, color=COLOR_ACCENT)
add_text(slide, Inches(0.75), Inches(6.45), Inches(12), Inches(0.4),
         "此问题已与导师和学长沟通确认。",
         font_size=12, color=COLOR_GRAY)

add_page_number(slide, 12, TOTAL_SLIDES)


# ==================== 第 13 页: 主要结果 — 9 数据集 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "03  全量实验结果 — 9 个数据集主指标",
              "35 / 35 指标全部提升，0 退化")

add_image(slide, FIG_DIR / "fig_main_results.png",
          Inches(0.5), Inches(1.1), width=Inches(12.3))

# 底部高亮
add_rect(slide, Inches(0.5), Inches(6.0), Inches(12.333), Inches(0.95),
         fill=RGBColor(0xDC, 0xFC, 0xE7), line=COLOR_GREEN, line_width=1)
add_text(slide, Inches(0.75), Inches(6.1), Inches(12), Inches(0.4),
         "🏆 核心成果：35 / 35 指标全部正向提升（无退化、无持平）",
         font_size=16, bold=True, color=COLOR_GREEN)
add_text(slide, Inches(0.75), Inches(6.5), Inches(12), Inches(0.4),
         "· FashionIQ / CIRCO / CIRR 使用默认参数 (β=0.7, α=0.9)    · GeneCIS 使用专用 prompt + task-adaptive 参数",
         font_size=11, color=COLOR_GRAY)

add_page_number(slide, 13, TOTAL_SLIDES)


# ==================== 第 14 页: 相对提升条形图 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "03  相对提升幅度对比",
              "GeneCIS 提升最惊人：+47% ~ +85%")

add_image(slide, FIG_DIR / "fig_relative.png",
          Inches(1.5), Inches(1.2), width=Inches(10.3))

# 右侧要点
add_text(slide, Inches(0.5), Inches(6.1), Inches(12), Inches(0.5),
         "▎ 重点解读", font_size=16, bold=True, color=COLOR_PRIMARY)

add_text(slide, Inches(0.5), Inches(6.55), Inches(6), Inches(0.4),
         "• FashionIQ / CIRCO / CIRR：+5% ~ +31% 相对提升，大规模图库仍有效",
         font_size=12, color=COLOR_GRAY)
add_text(slide, Inches(6.7), Inches(6.55), Inches(6), Inches(0.4),
         "• GeneCIS：+47% ~ +85% 相对提升，细粒度小图库场景尤其获益",
         font_size=12, color=COLOR_GRAY)

add_page_number(slide, 14, TOTAL_SLIDES)


# ==================== 第 15 页: FashionIQ 详细结果 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "03  FashionIQ 子集详细结果",
              "3 子集 × R@1/5/10/50 = 12 项指标全部提升")

add_image(slide, FIG_DIR / "fig_fashioniq.png",
          Inches(0.5), Inches(1.2), width=Inches(12.3))

# 底部数据
add_text(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(0.5),
         "▎ 关键数据点", font_size=16, bold=True, color=COLOR_PRIMARY)

points = [
    ("dress", "R@10: 15.80 → 19.29  (+3.49, +22.1%)",
     "R@50: 32.74 → 38.16  (+5.42)"),
    ("shirt", "R@10: 26.00 → 27.35  (+1.35, +5.2%)",
     "R@50: 42.94 → 44.84  (+1.90)"),
    ("toptee", "R@10: 23.09 → 27.35  (+4.26, +18.4%)",
     "R@50: 41.19 → 46.44  (+5.25)"),
]
for i, (sub, p1, p2) in enumerate(points):
    left = Inches(0.5 + i * 4.3)
    add_rect(slide, left, Inches(5.75), Inches(4.1), Inches(1.2),
             fill=COLOR_LIGHT, line=COLOR_PRIMARY, line_width=0.5)
    add_text(slide, left + Inches(0.2), Inches(5.85), Inches(3.7), Inches(0.4),
             sub, font_size=14, bold=True, color=COLOR_ACCENT)
    add_text(slide, left + Inches(0.2), Inches(6.3), Inches(3.7), Inches(0.3),
             p1, font_size=11, color=COLOR_GRAY, font_name=FONT_EN)
    add_text(slide, left + Inches(0.2), Inches(6.6), Inches(3.7), Inches(0.3),
             p2, font_size=11, color=COLOR_GRAY, font_name=FONT_EN)

add_page_number(slide, 15, TOTAL_SLIDES)


# ==================== 第 16 页: GeneCIS 专用方案 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "03  GeneCIS 专用方案",
              "短文本 + 小图库需要 task-adaptive 策略")

# 左侧：GeneCIS 特点
add_rect(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(3.1),
         fill=COLOR_LIGHT, line=COLOR_GRAY, line_width=0.5)
add_text(slide, Inches(0.75), Inches(1.4), Inches(5.8), Inches(0.4),
         "▎ GeneCIS 与其他数据集的差异", font_size=16,
         bold=True, color=COLOR_PRIMARY)

cmp_rows = [
    ("属性", "FIQ / CIRR / CIRCO", "GeneCIS"),
    ("修改文本", "句子 (15~30 词)", "1~2 词"),
    ("Gallery 大小", "几千 ~ 数十万", "约 14 张/查询"),
    ("任务特点", "全局语义匹配", "细粒度区分"),
]
base_top = Inches(2.0)
row_h = Inches(0.42)
widths = [Inches(1.8), Inches(2.3), Inches(1.9)]
col_xs = [Inches(0.75), Inches(2.65), Inches(5.0)]
for ri, row in enumerate(cmp_rows):
    top = base_top + row_h * ri
    bg = COLOR_PRIMARY if ri == 0 else (RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 else COLOR_LIGHT)
    add_rect(slide, Inches(0.7), top, Inches(5.9), row_h,
             fill=bg, line=RGBColor(0xE5, 0xE7, 0xEB), line_width=0.2)
    for ci, cell in enumerate(row):
        color = RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else COLOR_GRAY
        bold = (ri == 0) or (ci == 0)
        add_text(slide, col_xs[ci], top + Inches(0.07),
                 widths[ci], Inches(0.35),
                 cell, font_size=11, bold=bold, color=color)

# 右侧：专用方案
add_rect(slide, Inches(6.9), Inches(1.2), Inches(6), Inches(3.1),
         fill=COLOR_BOX_BG, line=COLOR_PRIMARY, line_width=0.8)
add_text(slide, Inches(7.15), Inches(1.4), Inches(5.6), Inches(0.4),
         "▎ 专用方案改动", font_size=16, bold=True, color=COLOR_PRIMARY)

changes = [
    ("1. 专用 prompt (V2)",
     "放宽「描述不超过修改文本」限制；要求输出具体视觉特征",
     "避免「modified/altered/changed」等模糊词"),
    ("2. 任务自适应 α/β",
     "每个子任务独立寻优，不用统一默认参数",
     ""),
]
top = Inches(2.0)
for (t, d1, d2) in changes:
    add_text(slide, Inches(7.15), top, Inches(5.6), Inches(0.4),
             t, font_size=13, bold=True, color=COLOR_ACCENT)
    add_text(slide, Inches(7.15), top + Inches(0.4), Inches(5.6), Inches(0.5),
             d1, font_size=11, color=COLOR_GRAY)
    if d2:
        add_text(slide, Inches(7.15), top + Inches(0.9), Inches(5.6), Inches(0.4),
                 d2, font_size=11, color=COLOR_GRAY)
    top += Inches(1.2)

# 底部结果表
add_text(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.4),
         "▎ GeneCIS 最终结果 (R@1)", font_size=16,
         bold=True, color=COLOR_PRIMARY)

genecis_rows = [
    ("子集", "Baseline", "三路融合", "绝对提升", "相对提升"),
    ("change_object",    "13.83", "25.51", "+11.68", "+84.5%"),
    ("focus_object",     "16.02", "23.62", "+7.60",  "+47.4%"),
    ("change_attribute", "12.65", "21.79", "+9.14",  "+72.3%"),
    ("focus_attribute",  "18.82", "27.83", "+9.01",  "+47.9%"),
]
col_ws = [Inches(3.0), Inches(2.1), Inches(2.1), Inches(2.3), Inches(2.3)]
col_xs = [Inches(0.6)]
for w in col_ws[:-1]:
    col_xs.append(col_xs[-1] + w)
base_top = Inches(5.0)
row_h = Inches(0.42)
for ri, row in enumerate(genecis_rows):
    top = base_top + row_h * ri
    bg = COLOR_PRIMARY if ri == 0 else (COLOR_LIGHT if ri % 2 else RGBColor(0xFF, 0xFF, 0xFF))
    tw = sum(col_ws, Emu(0))
    add_rect(slide, col_xs[0], top, tw, row_h,
             fill=bg, line=RGBColor(0xE5, 0xE7, 0xEB), line_width=0.2)
    for ci, cell in enumerate(row):
        color = RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else COLOR_GRAY
        if ri > 0 and ci == 3:
            color = COLOR_ACCENT
        bold = ri == 0 or ci == 0 or ci == 3
        align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        add_text(slide, col_xs[ci], top + Inches(0.07),
                 col_ws[ci], Inches(0.35),
                 cell, font_size=12, bold=bold, color=color, align=align,
                 font_name=FONT_EN if ci > 0 else FONT_CN)

add_page_number(slide, 16, TOTAL_SLIDES)


# ==================== 第 17 页: α/β 网格搜索热力图 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "04  α/β 网格搜索 — 不同任务偏好不同参数",
              "证明 task-adaptive 策略的价值")

add_image(slide, FIG_DIR / "fig_heatmap.png",
          Inches(0.5), Inches(1.1), width=Inches(12.3))

# 底部三列洞察
insights_3 = [
    ("change_object",
     "最优 β=0.30, α=0.95",
     "D₂ 占 70%，几乎不用代理图\n→ V7 精炼对物体变化最有效"),
    ("focus_object",
     "最优 β=1.00, α=0.90",
     "完全不用 D₂！V7 短描述限制\n让 D₂ 退化为 1-3 词，反而是噪声"),
    ("focus_attribute",
     "最优 β=0.60, α=0.85",
     "三路信号全发挥作用\n代理图贡献 15%"),
]
for i, (title, param, desc) in enumerate(insights_3):
    left = Inches(0.5 + i * 4.3)
    add_rect(slide, left, Inches(5.3), Inches(4.1), Inches(1.7),
             fill=COLOR_BOX_BG, line=COLOR_PRIMARY, line_width=0.5)
    add_text(slide, left + Inches(0.2), Inches(5.4), Inches(3.7), Inches(0.4),
             title, font_size=13, bold=True, color=COLOR_ACCENT)
    add_text(slide, left + Inches(0.2), Inches(5.8), Inches(3.7), Inches(0.4),
             param, font_size=11, bold=True,
             color=COLOR_PRIMARY, font_name=FONT_EN)
    add_text(slide, left + Inches(0.2), Inches(6.2), Inches(3.7), Inches(0.8),
             desc, font_size=11, color=COLOR_GRAY)

add_page_number(slide, 17, TOTAL_SLIDES)


# ==================== 第 18 页: 消融实验 — 各模块贡献 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "04  消融实验 — 各模块独立贡献",
              "三个模块各自解决不同问题，相互互补")

add_image(slide, FIG_DIR / "fig_ablation.png",
          Inches(0.5), Inches(1.15), width=Inches(12.3))

# 底部核心洞察
add_rect(slide, Inches(0.5), Inches(5.7), Inches(12.333), Inches(1.25),
         fill=COLOR_LIGHT, line=COLOR_PRIMARY, line_width=0.5)

add_text(slide, Inches(0.75), Inches(5.8), Inches(12), Inches(0.4),
         "🔑 核心洞察：单一模块的最优配置在不同数据集上截然相反",
         font_size=14, bold=True, color=COLOR_PRIMARY)
add_text(slide, Inches(0.75), Inches(6.25), Inches(12), Inches(0.3),
         "• 在 CIRR 上：V7 精炼是主力（+6.5），Ensemble 反而回调一部分",
         font_size=11, color=COLOR_GRAY)
add_text(slide, Inches(0.75), Inches(6.55), Inches(12), Inches(0.3),
         "• 在 FashionIQ dress 上：V7 单独使用会退化（→13），必须 Ensemble 兜底",
         font_size=11, color=COLOR_GRAY)

add_page_number(slide, 18, TOTAL_SLIDES)


# ==================== 第 19 页: 工作量 & 工程统计 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "04  工作量与工程统计",
              "完整实验链路 + 复现 + 创新")

stats = [
    ("9", "全量数据集评估", COLOR_PRIMARY),
    ("35/35", "指标全部提升", COLOR_GREEN),
    ("6", "Prompt 版本迭代", COLOR_ACCENT),
    ("56+", "网格搜索参数组合", COLOR_PRIMARY),
    ("2.2w+", "MLLM API 调用", COLOR_GREEN),
    ("~1.9w", "代理图生成", COLOR_ACCENT),
]

for i, (num, desc, c) in enumerate(stats):
    col = i % 3
    row = i // 3
    left = Inches(0.7 + col * 4.2)
    top = Inches(1.25 + row * 1.8)

    add_rect(slide, left, top, Inches(3.9), Inches(1.6),
             fill=COLOR_LIGHT, line=c, line_width=1.5)
    add_text(slide, left, top + Inches(0.15), Inches(3.9), Inches(0.8),
             num, font_size=44, bold=True, color=c,
             align=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_text(slide, left, top + Inches(1.05), Inches(3.9), Inches(0.4),
             desc, font_size=13, color=COLOR_GRAY,
             align=PP_ALIGN.CENTER)

# 底部实现统计
add_rect(slide, Inches(0.5), Inches(5.0), Inches(12.333), Inches(1.95),
         fill=COLOR_BOX_BG, line=COLOR_PRIMARY, line_width=0.5)
add_text(slide, Inches(0.75), Inches(5.15), Inches(12), Inches(0.4),
         "▎ 实现细节",
         font_size=16, bold=True, color=COLOR_PRIMARY)

impl = [
    ("环境", "Linux (GPU + MLLM API) + Windows (RTX 4060 + CLIP 编码)"),
    ("API 花费", "约 140 元 (MLLM ≈90 元 + 文生图 ≈50 元)"),
    ("代码规模", "src + scripts 共约 5000 行 Python；论文 LaTeX 源 ≈26 页 PDF"),
    ("GitHub", "https://github.com/Haookok/osrcir-threeway-fusion"),
]
for i, (k, v) in enumerate(impl):
    top = Inches(5.6 + i * 0.32)
    add_text(slide, Inches(0.85), top, Inches(2.0), Inches(0.3),
             f"· {k}", font_size=12, bold=True, color=COLOR_PRIMARY)
    add_text(slide, Inches(2.95), top, Inches(9.6), Inches(0.3),
             v, font_size=11, color=COLOR_GRAY)

add_page_number(slide, 19, TOTAL_SLIDES)


# ==================== 第 20 页: 方法优势与局限 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "05  方法优势与局限",
              "客观评价：做到了什么 & 还有哪些不足")

# 左侧：优势
add_rect(slide, Inches(0.5), Inches(1.15), Inches(6.2), Inches(5.7),
         fill=RGBColor(0xDC, 0xFC, 0xE7), line=COLOR_GREEN, line_width=1)
add_text(slide, Inches(0.75), Inches(1.3), Inches(5.8), Inches(0.4),
         "✓ 方法优势", font_size=18, bold=True, color=COLOR_GREEN)

pros = [
    ("无需训练",
     "保持零样本方法的通用性，易于跨数据集部署"),
    ("视觉校验闭环",
     "首次在纯文本 ZS-CIR 流程中引入图像空间信号"),
    ("稳健性保证",
     "Ensemble 机制让整体不因精炼失误而退化"),
    ("可扩展架构",
     "可调 α/β 和 prompt，容纳多种任务类型"),
    ("可复现性强",
     "全部代码、数据、日志已开源"),
]
for i, (t, d) in enumerate(pros):
    top = Inches(1.9 + i * 0.95)
    add_text(slide, Inches(0.85), top, Inches(5.6), Inches(0.4),
             f"• {t}", font_size=14, bold=True, color=COLOR_PRIMARY)
    add_text(slide, Inches(1.05), top + Inches(0.4), Inches(5.4), Inches(0.5),
             d, font_size=12, color=COLOR_GRAY)

# 右侧：局限
add_rect(slide, Inches(6.9), Inches(1.15), Inches(6), Inches(5.7),
         fill=RGBColor(0xFE, 0xE2, 0xE2), line=COLOR_ACCENT, line_width=1)
add_text(slide, Inches(7.15), Inches(1.3), Inches(5.6), Inches(0.4),
         "△ 方法局限", font_size=18, bold=True, color=COLOR_ACCENT)

cons = [
    ("推理成本增加",
     "多 1 次文生图 + 1 次 MLLM 精炼调用"),
    ("参数敏感性",
     "统一默认参数并非所有场景最优"),
    ("依赖 MLLM 能力上限",
     "本工作因 API 限制使用 Qwen-VL-Max 而非 GPT-4o"),
    ("代理图噪声",
     "文生图模型自身的幻觉仍可能影响辅助信号"),
    ("评测 split 限制",
     "CIRR/CIRCO 使用 val split（OSrCIR test 不公开）"),
]
for i, (t, d) in enumerate(cons):
    top = Inches(1.9 + i * 0.95)
    add_text(slide, Inches(7.25), top, Inches(5.4), Inches(0.4),
             f"• {t}", font_size=14, bold=True, color=COLOR_PRIMARY)
    add_text(slide, Inches(7.45), top + Inches(0.4), Inches(5.2), Inches(0.5),
             d, font_size=12, color=COLOR_GRAY)

add_page_number(slide, 20, TOTAL_SLIDES)


# ==================== 第 21 页: 总结与展望 ====================
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, "05  总结与展望", "Conclusion & Future Work")

# 左：主要贡献
add_rect(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.7),
         fill=COLOR_BOX_BG, line=COLOR_PRIMARY, line_width=1)
add_text(slide, Inches(0.75), Inches(1.35), Inches(5.8), Inches(0.4),
         "▎ 主要贡献", font_size=18, bold=True, color=COLOR_PRIMARY)

contribs = [
    "首次在 ZS-CIR 中提出视觉代理机制，\n为纯文本检索流程引入图像空间信号",
    "设计 V7 反幻觉提示词，让 AI 生成图像\n作为诊断工具而非描述来源",
    "提出 CLIP 特征空间的描述融合策略，\n同时提升性能并保持稳健性",
    "在 9 个标准数据集上实现 35/35 指标\n全部正向提升（无退化、无持平）",
]
for i, c in enumerate(contribs):
    top = Inches(1.95 + i * 1.2)
    add_text(slide, Inches(0.85), top, Inches(0.4), Inches(0.4),
             f"①②③④"[i:i+1], font_size=24, bold=True,
             color=COLOR_ACCENT, font_name=FONT_EN)
    add_text(slide, Inches(1.3), top + Inches(0.05), Inches(5.2), Inches(1.1),
             c, font_size=12, color=COLOR_GRAY)

# 右：未来工作
add_rect(slide, Inches(6.9), Inches(1.2), Inches(6), Inches(5.7),
         fill=COLOR_LIGHT, line=COLOR_GRAY, line_width=0.5)
add_text(slide, Inches(7.15), Inches(1.35), Inches(5.6), Inches(0.4),
         "▎ 未来工作", font_size=18, bold=True, color=COLOR_PRIMARY)

future = [
    ("查询自适应参数",
     "学习预测器根据 query 自动选择 α/β"),
    ("更强 MLLM",
     "接入 GPT-4o/Gemini/Claude 提升 D₁ 上限"),
    ("代理图优化",
     "多代理图集成、SDXL/Flux 替换、代理图质量过滤"),
    ("扩展到有监督设定",
     "验证框架在有监督 CIR 上的适用性"),
]
for i, (t, d) in enumerate(future):
    top = Inches(1.95 + i * 1.2)
    add_text(slide, Inches(7.25), top, Inches(0.4), Inches(0.4),
             f"①②③④"[i:i+1], font_size=24, bold=True,
             color=COLOR_GREEN, font_name=FONT_EN)
    add_text(slide, Inches(7.7), top + Inches(0.05), Inches(4.9), Inches(0.4),
             t, font_size=13, bold=True, color=COLOR_PRIMARY)
    add_text(slide, Inches(7.7), top + Inches(0.5), Inches(4.9), Inches(0.7),
             d, font_size=11, color=COLOR_GRAY)

add_page_number(slide, 21, TOTAL_SLIDES)


# ==================== 第 22 页: 致谢 ====================
slide = prs.slides.add_slide(BLANK)

# 背景渐变效果（用两层矩形模拟）
bg1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0), Inches(0),
                              SLIDE_W, SLIDE_H)
bg1.fill.solid()
bg1.fill.fore_color.rgb = COLOR_PRIMARY
bg1.line.fill.background()
bg1.shadow.inherit = False

# 装饰线
dec_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(4.5), Inches(3.3),
                                   Inches(4.3), Inches(0.05))
dec_line.fill.solid()
dec_line.fill.fore_color.rgb = COLOR_ACCENT
dec_line.line.fill.background()
dec_line.shadow.inherit = False

# 主文字
add_text(slide, Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.0),
         "谢  谢", font_size=80, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.5), Inches(12.333), Inches(0.6),
         "Thanks for Your Attention",
         font_size=22, color=RGBColor(0xCC, 0xDB, 0xFE),
         align=PP_ALIGN.CENTER, font_name=FONT_EN)

add_text(slide, Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5),
         "恳请各位老师批评指正",
         font_size=20, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(5.2), Inches(12.333), Inches(0.5),
         "Questions & Answers",
         font_size=16, color=RGBColor(0xCC, 0xDB, 0xFE),
         align=PP_ALIGN.CENTER, font_name=FONT_EN)

add_text(slide, Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.4),
         "杨昊明  |  2026 年 5 月  |  同济大学毕业设计（论文）答辩",
         font_size=12, color=RGBColor(0xCC, 0xDB, 0xFE),
         align=PP_ALIGN.CENTER)


# ---------------- 保存 ----------------
OUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT_PPTX))
print(f"✓ PPT 已生成: {OUT_PPTX}")
print(f"  共 {len(prs.slides)} 页")
