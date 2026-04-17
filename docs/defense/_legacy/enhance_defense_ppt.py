# -*- coding: utf-8 -*-
"""
Enhance the defense PPT:
  1) add Speaker Notes (presenter script) to every slide
  2) insert a "Core Numbers Cheat-Sheet" slide before the Thank-You page

Run:
    python enhance_defense_ppt.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = os.path.dirname(os.path.abspath(__file__))
PPT_PATH = os.path.join(BASE, "毕业答辩_OSrCIR三路融合.pptx")

TJ_PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
DEEP_BLUE = RGBColor(0x14, 0x2B, 0x5C)
ACCENT = RGBColor(0xE0, 0x4B, 0x3F)
LIGHT_BG = RGBColor(0xF5, 0xF3, 0xF9)
GRAY_TXT = RGBColor(0x55, 0x55, 0x55)
DARK_TXT = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ----- Speaker notes for each of the 17 original slides -----
SPEAKER_NOTES = [
    # 0: 封面
    "[开场 30 秒]\n各位老师好，我是同济大学 2026 届本科毕业生杨浩铭，我的毕业设计题目是《基于视觉代理与描述融合的零样本组合式图像检索改进》，指导老师是 X 老师。本课题复现并改进了 CVPR 2025 Highlight 论文 OSrCIR。下面我将用大约 12 分钟向各位老师汇报。",
    # 1: 目录
    "[目录 15 秒]\n本次汇报分为六部分：研究背景与任务定义、基线方法瓶颈、三路融合方法、实验结果、消融分析、总结与展望。",
    # 2: 背景与任务
    "[背景 60 秒]\n组合式图像检索 CIR 接受 参考图 + 修改文本 双模态 query，在大图库中检索满足 修改意图 的目标图。它比纯图搜图更灵活，比 VQA 更面向检索，贴近电商、内容创作等真实场景。\n零样本 ZS-CIR 要求不允许在 CIR 数据集上训练或微调，只能靠预训练好的 CLIP、MLLM 直接推理，省去昂贵的三元组标注，同时在新领域零成本泛化。\nCVPR 2025 的 OSrCIR 用 MLLM + Chain-of-Thought 直接生成目标图描述，是当前 SOTA 之一，也是我改进的基线。",
    # 3: 瓶颈
    "[基线瓶颈 60 秒]\n我在复现过程中发现三个痛点：\n第一，OSrCIR 全程在文本空间做检索，缺少图像级 anchor，MLLM 描述一旦偏差就没有纠正途径；\n第二，MLLM 偶尔会对目标的空间布局、颜色深浅理解错位，因为它看不到目标的可视化形态；\n第三，单一描述鲁棒性弱，跨数据集表现不稳定。\n围绕这三点，我设计了三路融合方案。",
    # 4: 方法总览
    "[方法总览 75 秒]\n三路融合可以一句话概括：\nscore 等于 α 倍的 文本与图库相似度 加上 (1-α) 倍的 代理图与图库相似度 ；文本特征是 β 倍 CLIP(D1) 加 (1-β) 倍 CLIP(D2) 的加权并归一化。\n第一路：MLLM 思考参考图+修改文本，生成初始描述 D1；\n第二路：把 D1 送给 MiniMax 文生图，得到代理图 proxy，再让 MLLM 看着 (参考图 + 代理图 + 修改文本) 用 V7 反幻觉 prompt 输出精炼描述 D2；\n第三路：代理图本身也参与检索。\n默认 α=0.9、β=0.7 来自 200 样本网格搜索在 9 个数据集上的平均最稳组合。",
    # 5: Visual Proxy
    "[创新点 1 · Visual Proxy 50 秒]\n核心动机：让 MLLM 的描述 能被看见 。\n做法：把 D1 送进 MiniMax image-01，得到代理图；用 CLIP 编码，参与检索打分。它不需要像真实目标那样精确，只需要 描述同一个语义 即可——因为 CLIP 的相似度是语义级的。\n消融数据：仅加代理图这一路，FashionIQ dress R@10 从 18.3 升到 19.3，确认有效。",
    # 6: V7 反幻觉
    "[创新点 2 · V7 反幻觉 75 秒]\n这是我花时间最多的一步。最初的 prompt 让 MLLM 对比参考图和代理图 综合写新描述 ，结果 CIRR R@10 反而从 67 跌到 64.5，相当于负收益。\n逐样本分析发现：MLLM 把代理图里虚构的背景、材质、光影全写进了描述，长度膨胀到原来的 1.9 倍——这就是幻觉。\nV7 里我写了三条硬性指令：\n1) 明确声明代理图含幻觉细节、不得引用；\n2) 代理图只作为 诊断工具 ，用来判断 修改是否已经发生 ；\n3) 输出必须是短描述、长度不超过修改文本。\nV7 把 CIRR 从 67 拉到 73.5，提升 6.5 个点。",
    # 7: Description Ensemble
    "[创新点 3 · Description Ensemble 50 秒]\nV7 在 CIRR 表现极佳，但在 FashionIQ shirt 会丢失关键词、退化 4 个点。原因是 完全用 D2 替换 D1 粒度太粗。\n解决：特征级加权融合 f_text = normalize( β · CLIP(D1) + (1-β) · CLIP(D2) )，β=0.7。\n直观上，70% 保留原始描述 兜底 ，30% 让精炼描述做 增量纠偏 。\n结果：牺牲了 CIRR 的峰值，但换来在 9 个数据集上的一致正向。",
    # 8: 实验设置
    "[实验设置 30 秒]\nCLIP 用 ViT-L/14、MLLM 用 Qwen-VL-Max、文生图用 MiniMax；9 个数据集：FashionIQ 三子集 + CIRCO + CIRR + GeneCIS 四子集；所有数据集 baseline 和 3-way 用相同 MLLM、相同 CLIP、相同 gallery；GeneCIS 额外做 per-subset 网格搜索得到 task-adaptive 最优 α/β。",
    # 9: 主结果
    "[主结果 90 秒]\n这张图是 9 个数据集 baseline vs 三路融合对比。结论很清楚——35 个指标全部正向提升，没有一个持平、没有一个下降。\n重点几个：CIRCO mAP@10 从 16.21 到 21.26，相对 +31%；CIRR R@1 从 22.96 到 25.90，+13%；GeneCIS change_object R@1 从 13.83 到 25.51，相对 +85%，全场最大。\n特别说明：我没直接对比原论文 GPT-4o 数字，因为那是 MLLM 差异而非方法差异，baseline 已经用相同 Qwen 复现过。",
    # 10: FashionIQ 细化
    "[FashionIQ 细化 45 秒]\nFashionIQ 分 dress/shirt/toptee 三个子集。dress R@10 +3.49 (+22%)、toptee +4.26 (+18%)；shirt 提升仅 +5% 但仍正向——这也解释了为什么需要 ensemble——V7 在 shirt 上会退化，ensemble 让它 不退化且仍然有提升 。",
    # 11: α/β 网格搜索
    "[α/β 网格搜索 60 秒]\n这张热力图是 GeneCIS 子集 R@1 随 α/β 变化。观察：默认 α=0.9 β=0.7 在绝大多数子集 top-3；不同子集最优差异显著，change_object 偏好 α=0.95 β=0.50，focus_object 反而是 α=0.80 β=1.00——说明 task-adaptive 是一个有价值的研究方向。",
    # 12: 组件消融
    "[组件消融 45 秒]\n这张条形图把三路融合拆成三个独立组件：只加 ensemble、只加 V7、只加 proxy。三个组件都有独立贡献，完整三路融合取得最大增益——证明三者互补而非冗余。proxy + ensemble 组合同时给了图像 anchor 和描述鲁棒性。",
    # 13: 相对提升
    "[相对提升概览 30 秒]\n按相对提升百分比排序。最低 FashionIQ shirt +5.2%，最高 GeneCIS change_object +84.5%；所有 bar 全部在 0 线以上——这是我最引以为豪的一张图。",
    # 14: 工作贡献
    "[工作贡献 45 秒]\n三点贡献：\n1) 方法论：首次把 Visual Proxy 引入 ZS-CIR，让检索同时获得文本和图像双视角；\n2) 工程：提出 V7 反幻觉 prompt 原则，系统性解决 AI 生成输入回流导致幻觉 ；\n3) 鲁棒性：提出 Description Ensemble，把 V7 的局部峰值转化为 9 数据集一致正向。",
    # 15: 局限与未来
    "[局限与未来 45 秒]\n不足：对 MLLM 品牌有依赖、proxy 生成每张 2-3 秒不适合实时、缺人工评估。\n未来：端到端联合训练、换 SigLIP / EVA-CLIP、引入人工评估、扩展到视频 CIR。",
    # 16: 致谢
    "[致谢 15 秒]\n感谢指导老师 X 老师、感谢同组同学、感谢各位答辩评审老师。汇报完毕，欢迎提问。",
]

CHEAT_NOTES = "[备用 slide · Q&A 可切到此页]\n核心数字速查。老师问到具体指标时指向此页精确回答：数据集 9；查询 ~18K；最大 gallery CIRCO 123K；默认 α=0.9 β=0.7；全部 35 指标正向；prompt 迭代 4 轮；API 调用约 15 万次；总成本约 140 元。"


def _add_text(slide, left, top, width, height, text, size=18, bold=False,
              color=DARK_TXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              font_name="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def _row(slide, x, y, w, h, label, value, value_color):
    _add_text(slide, x, y, Inches(w * 0.55), Inches(h), label,
              size=15, color=DARK_TXT, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, x + Inches(w * 0.55), y, Inches(w * 0.45), Inches(h),
              value, size=16, bold=True, color=value_color,
              anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


def build_cheatsheet(prs):
    sw, sh = prs.slide_width, prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.shadow.inherit = False

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Inches(0.9))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = TJ_PURPLE
    band.shadow.inherit = False

    _add_text(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
              "核心数字速查  ·  Answer Cheat-Sheet",
              size=26, bold=True, color=WHITE)

    left_x = Inches(0.7)
    top_y = Inches(1.25)
    _add_text(slide, left_x, top_y, Inches(6), Inches(0.4),
              "▎ 任务与设置", size=17, bold=True, color=DEEP_BLUE)

    rows_left = [
        ("数据集总数", "9 个"),
        ("FashionIQ 子集", "dress / shirt / toptee"),
        ("CIRCO / CIRR", "220 / 4181 queries"),
        ("GeneCIS 子集", "ch/fc × obj/attr"),
        ("查询总数", "~ 18 K"),
        ("最大 gallery", "CIRCO 123 K"),
        ("CLIP backbone", "ViT-L/14"),
        ("MLLM", "Qwen-VL-Max"),
        ("T2I", "MiniMax image-01"),
        ("默认 α / β", "0.90 / 0.70"),
    ]
    y = top_y + Inches(0.45)
    for lb, vl in rows_left:
        _row(slide, left_x, y, 6.0, 0.35, lb, vl, DEEP_BLUE)
        y += Inches(0.38)

    right_x = Inches(7.0)
    _add_text(slide, right_x, top_y, Inches(6), Inches(0.4),
              "▎ 主要结果 (baseline → 3-way)", size=17, bold=True, color=DEEP_BLUE)

    rows_right = [
        ("FIQ dress R@10", "15.80 → 19.29 (+22%)"),
        ("FIQ shirt R@10", "26.00 → 27.35 (+5%)"),
        ("FIQ toptee R@10", "23.09 → 27.35 (+18%)"),
        ("CIRCO mAP@10", "16.21 → 21.26 (+31%)"),
        ("CIRR R@1", "22.96 → 25.90 (+13%)"),
        ("GeneCIS ch_obj R@1", "13.83 → 25.51 (+85%)"),
        ("GeneCIS fc_obj R@1", "16.02 → 23.62 (+47%)"),
        ("GeneCIS ch_attr R@1", "12.65 → 21.79 (+72%)"),
        ("GeneCIS fc_attr R@1", "18.82 → 27.83 (+48%)"),
        ("总指标", "35 / 35 正向"),
    ]
    y = top_y + Inches(0.45)
    for lb, vl in rows_right:
        _row(slide, right_x, y, 6.0, 0.35, lb, vl, ACCENT)
        y += Inches(0.38)

    _add_text(slide, Inches(0.7), Inches(6.45), Inches(12), Inches(0.45),
              "规模 · prompt 迭代 4 轮  |  API 调用 ~ 15 万次  |  总成本 ≈ 140 元",
              size=13, color=GRAY_TXT, align=PP_ALIGN.CENTER)
    return slide


def move_before_last(prs):
    xml = prs.slides._sldIdLst
    ids = list(xml)
    new_one = ids[-1]
    thanks = ids[-2]
    xml.remove(new_one)
    xml.insert(list(xml).index(thanks), new_one)


def write_notes(prs):
    for i, slide in enumerate(prs.slides):
        if i < len(SPEAKER_NOTES):
            note = SPEAKER_NOTES[i]
        else:
            note = CHEAT_NOTES
        tf = slide.notes_slide.notes_text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = note
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(11)


def main():
    prs = Presentation(PPT_PATH)
    orig = len(prs.slides)
    print("[INFO] 打开 PPT 共 %d 页" % orig)

    cheat = build_cheatsheet(prs)
    print("[INFO] 已添加 核心数字速查 slide")
    move_before_last(prs)
    print("[INFO] 已移至 致谢 页之前")

    write_notes(prs)
    print("[INFO] 已为全部 %d 页写入 Speaker Notes" % len(prs.slides))

    prs.save(PPT_PATH)
    print("[DONE] 保存至:", PPT_PATH)


if __name__ == "__main__":
    main()
