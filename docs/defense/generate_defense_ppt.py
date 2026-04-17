"""
生成毕业答辩 PPT

运行:
    python generate_defense_ppt.py

输出:
    /home/haomingyang03/code/osrcir/docs/defense/毕业答辩_OSrCIR三路融合.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---------- 配色 (同济紫 + 深蓝 + 白) ----------
TJ_PURPLE = RGBColor(0x6A, 0x1B, 0x9A)   # 同济主色
DEEP_BLUE = RGBColor(0x14, 0x2B, 0x5C)   # 深蓝
ACCENT = RGBColor(0xE0, 0x4B, 0x3F)      # 红色强调
LIGHT_BG = RGBColor(0xF5, 0xF3, 0xF9)    # 浅紫背景
GRAY_TXT = RGBColor(0x55, 0x55, 0x55)
DARK_TXT = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xEF, 0x6C, 0x00)

# ---------- 文件路径 ----------
BASE = "/home/haomingyang03/code/osrcir/docs"
FIG = os.path.join(BASE, "thesis/figures")
OUT = os.path.join(BASE, "defense/毕业答辩_OSrCIR三路融合.pptx")

# ---------- 幻灯片尺寸: 16:9 ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# ---------- 通用工具 ----------
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # 置底
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_w is not None:
            shape.line.width = line_w
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_TXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_image_fit(slide, path, cx, cy, max_w, max_h):
    """居中放图片，按比例fit到max_w/max_h框内。"""
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w = int(iw * ratio)
    h = int(ih * ratio)
    x = int(cx - w / 2)
    y = int(cy - h / 2)
    return slide.shapes.add_picture(path, x, y, width=w, height=h)


def page_header(slide, title, subtitle=None, page=None, total=None):
    """标题栏 (顶部色块 + 标题文字 + 页码)"""
    # 顶部色条
    add_rect(slide, 0, 0, SW, Inches(0.9), fill=TJ_PURPLE)
    # 左侧装饰竖条
    add_rect(slide, 0, Inches(0.9), Inches(0.08), SH - Inches(0.9), fill=TJ_PURPLE)
    # 标题
    add_text(slide, Inches(0.4), Inches(0.12), Inches(10), Inches(0.5),
             title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.52), Inches(10), Inches(0.3),
                 subtitle, size=12, color=RGBColor(0xDD, 0xDD, 0xEE))
    # 页码
    if page and total:
        add_text(slide, Inches(11.8), Inches(0.3), Inches(1.4), Inches(0.4),
                 f"{page} / {total}", size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT)


def footer(slide, text="同济大学 · 2026 届本科毕业设计答辩 · 杨浩铭"):
    add_text(slide, Inches(0.4), SH - Inches(0.4),
             Inches(12.5), Inches(0.3),
             text, size=10, color=GRAY_TXT, align=PP_ALIGN.CENTER)


# ==================================================
# 封面 (slide 1)
# ==================================================
def slide_cover():
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, LIGHT_BG)
    # 大色块
    add_rect(slide, 0, 0, SW, Inches(3.2), fill=TJ_PURPLE)
    add_rect(slide, 0, Inches(3.2), SW, Inches(0.08), fill=ACCENT)

    add_text(slide, Inches(0.8), Inches(0.6), Inches(12), Inches(0.5),
             "同济大学 2026 届本科毕业设计答辩", size=18, color=WHITE,
             align=PP_ALIGN.LEFT)

    add_text(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(1.4),
             "基于视觉代理与描述融合的", size=40, bold=True, color=WHITE)
    add_text(slide, Inches(0.8), Inches(1.9), Inches(12), Inches(1.4),
             "零样本组合式图像检索改进", size=40, bold=True, color=WHITE)

    add_text(slide, Inches(0.8), Inches(2.8), Inches(12), Inches(0.4),
             "— 对 CVPR 2025 Highlight OSrCIR 的三路融合改进 —",
             size=18, color=RGBColor(0xE0, 0xD0, 0xF0))

    # 信息块
    info_y = Inches(4.0)
    add_text(slide, Inches(0.8), info_y, Inches(2.5), Inches(0.5),
             "答辩人", size=16, color=TJ_PURPLE, bold=True)
    add_text(slide, Inches(0.8), info_y + Inches(0.4), Inches(3), Inches(0.5),
             "杨 浩 铭", size=22, bold=True, color=DARK_TXT)

    add_text(slide, Inches(4.5), info_y, Inches(3), Inches(0.5),
             "指导教师", size=16, color=TJ_PURPLE, bold=True)
    add_text(slide, Inches(4.5), info_y + Inches(0.4), Inches(4), Inches(0.5),
             "XXX  教授", size=22, bold=True, color=DARK_TXT)

    add_text(slide, Inches(8.5), info_y, Inches(3), Inches(0.5),
             "专业 / 学院", size=16, color=TJ_PURPLE, bold=True)
    add_text(slide, Inches(8.5), info_y + Inches(0.4), Inches(4.5), Inches(0.5),
             "计算机科学与技术 · 电子与信息工程学院",
             size=16, bold=True, color=DARK_TXT)

    # 底部基础信息
    add_rect(slide, 0, Inches(6.7), SW, Inches(0.8), fill=DEEP_BLUE)
    add_text(slide, Inches(0.8), Inches(6.85), Inches(12), Inches(0.5),
             "关键词:  零样本组合式图像检索 · 多模态大模型 · 视觉代理 · 描述融合 · CLIP",
             size=14, color=WHITE, align=PP_ALIGN.CENTER)


# ==================================================
# 目录 (slide 2)
# ==================================================
def slide_toc(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "目  录", "Contents", page=2, total=total)

    items = [
        ("01", "研究背景与意义", "零样本组合式图像检索 (ZS-CIR) 与 OSrCIR 基线"),
        ("02", "研究内容与问题", "基线瓶颈：描述缺乏视觉验证 / MLLM 易幻觉"),
        ("03", "方法：三路融合", "Visual Proxy + V7 反幻觉 CoT + Description Ensemble"),
        ("04", "实验与结果", "9 个标准数据集全量实验，35/35 指标全部提升"),
        ("05", "消融与分析", "α/β 网格搜索、prompt 迭代、跨数据集泛化"),
        ("06", "总结与展望", "工作贡献、局限性、后续可拓展方向"),
    ]
    y0 = Inches(1.4)
    row_h = Inches(0.85)
    for i, (idx, title, sub) in enumerate(items):
        y = y0 + row_h * i
        # 编号圆块
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y,
                                       Inches(0.8), Inches(0.8))
        circ.fill.solid()
        circ.fill.fore_color.rgb = TJ_PURPLE
        circ.line.fill.background()
        circ.shadow.inherit = False
        tf = circ.text_frame
        tf.text = idx
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(20)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.name = "Microsoft YaHei"

        # 标题
        add_text(slide, Inches(2.3), y + Inches(0.05), Inches(5),
                 Inches(0.4), title, size=22, bold=True, color=DARK_TXT)
        # 副标题
        add_text(slide, Inches(2.3), y + Inches(0.45), Inches(9),
                 Inches(0.3), sub, size=13, color=GRAY_TXT)

    footer(slide)


# ==================================================
# 01 研究背景
# ==================================================
def slide_background(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "01  研究背景与任务定义",
                "Zero-Shot Composed Image Retrieval (ZS-CIR)",
                page=3, total=total)

    # 左侧定义
    add_rect(slide, Inches(0.45), Inches(1.2), Inches(6.2), Inches(5.8),
             fill=LIGHT_BG)
    add_text(slide, Inches(0.7), Inches(1.3), Inches(6), Inches(0.5),
             "任务定义", size=20, bold=True, color=TJ_PURPLE)
    add_text(slide, Inches(0.7), Inches(1.8), Inches(6), Inches(1.8),
             "给定 (参考图, 修改文本) 二元组\n"
             "  —— 修改文本描述: “把裙子换成黑色短款”\n"
             "在大图库 (万~十万级) 中检索符合描述的目标图像。\n"
             "\n约束: 零样本 (Zero-Shot)\n"
             "  —— 不允许在下游数据集做任何训练或微调",
             size=14, color=DARK_TXT)

    add_text(slide, Inches(0.7), Inches(4.05), Inches(6), Inches(0.5),
             "应用场景", size=20, bold=True, color=TJ_PURPLE)
    bullets = [
        "• 电商服装检索: 以图+文细化筛选 (Amazon, 阿里)",
        "• 内容创作素材库: 基于已有图修改意图找素材",
        "• 安防/医疗: 带条件修改的相似案例匹配",
        "• 通用基础模型 (CLIP / MLLM) 下游通用能力评测",
    ]
    add_text(slide, Inches(0.7), Inches(4.5), Inches(6), Inches(2.4),
             "\n".join(bullets), size=13, color=DARK_TXT)

    # 右侧: Baseline 论文 OSrCIR
    add_rect(slide, Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.8),
             fill=RGBColor(0xEC, 0xEE, 0xF5))
    add_text(slide, Inches(7.1), Inches(1.3), Inches(6), Inches(0.5),
             "基线论文: OSrCIR (CVPR 2025 Highlight)",
             size=18, bold=True, color=DEEP_BLUE)
    add_text(slide, Inches(7.1), Inches(1.8), Inches(5.7), Inches(0.4),
             "arXiv 2412.11077 · 2024.12 · One-Stage reflective CIR",
             size=12, color=GRAY_TXT)

    add_text(slide, Inches(7.1), Inches(2.3), Inches(5.7), Inches(0.45),
             "核心流程 (单轮 CoT 推理)", size=15, bold=True, color=TJ_PURPLE)
    pipeline_lines = [
        "① 参考图  +  修改文本",
        "          ↓  (MLLM with CoT prompt)",
        "② 目标描述 D  (纯文本)",
        "          ↓  (CLIP ViT-L/14)",
        "③ 文本特征   →   与图库检索",
    ]
    add_text(slide, Inches(7.3), Inches(2.8), Inches(5.5), Inches(2.2),
             "\n".join(pipeline_lines), size=14, color=DARK_TXT)

    add_text(slide, Inches(7.1), Inches(5.2), Inches(5.7), Inches(0.4),
             "主要贡献", size=15, bold=True, color=TJ_PURPLE)
    add_text(slide, Inches(7.3), Inches(5.55), Inches(5.7), Inches(1.4),
             "• 将 MLLM 作为“反射式”理解器引入 ZS-CIR\n"
             "• 零训练,  纯 prompt + CLIP 对齐即可刷新 SOTA\n"
             "• 统一覆盖 FashionIQ / CIRCO / CIRR / GeneCIS",
             size=13, color=DARK_TXT)
    footer(slide)


# ==================================================
# 02 研究内容 / 痛点
# ==================================================
def slide_problem(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "02  基线瓶颈与改进切入点",
                "单轮推理的结构性缺陷 · MLLM 幻觉风险",
                page=4, total=total)

    # 三列卡片
    cards = [
        ("1", "缺乏视觉验证信号",
         "MLLM 仅凭参考图 + 文本做一次推理输出 D1,\n"
         "没有图像侧 (image space) 的反馈闭环。\n"
         "一旦遗漏关键修改, 错误会直接传导到检索。"),
        ("2", "MLLM 幻觉风险高",
         "直接在 prompt 中让 MLLM 自由精炼描述,\n"
         "容易引入代理图中的 AI 虚构细节 (背景/材质)\n"
         "导致描述膨胀 ~ 1.9×, CLIP 空间反而远离目标。"),
        ("3", "跨任务策略不统一",
         "FashionIQ / CIRCO / CIRR / GeneCIS 任务特性\n"
         "差异大 (时尚属性 / 开放场景 / 细粒度改动),\n"
         "单一 prompt 策略难以在 9 个子集全面正向。"),
    ]
    card_y = Inches(1.35)
    card_h = Inches(3.4)
    card_w = Inches(4.0)
    gap = Inches(0.28)
    x0 = Inches(0.45)
    for i, (n, title, body) in enumerate(cards):
        x = x0 + (card_w + gap) * i
        # 阴影底块
        add_rect(slide, x + Inches(0.08), card_y + Inches(0.08),
                 card_w, card_h, fill=RGBColor(0xDD, 0xDD, 0xE8))
        add_rect(slide, x, card_y, card_w, card_h, fill=WHITE,
                 line=TJ_PURPLE, line_w=Pt(1.2))
        # 顶部色条
        add_rect(slide, x, card_y, card_w, Inches(0.55), fill=TJ_PURPLE)
        add_text(slide, x + Inches(0.25), card_y + Inches(0.08),
                 Inches(3.6), Inches(0.4),
                 f"难点 {n}  ·  {title}", size=15, bold=True, color=WHITE)
        add_text(slide, x + Inches(0.25), card_y + Inches(0.75),
                 card_w - Inches(0.5), card_h - Inches(0.9),
                 body, size=13, color=DARK_TXT)

    # 下方: 改进思路
    y2 = Inches(5.0)
    add_rect(slide, Inches(0.45), y2, Inches(12.4), Inches(1.95),
             fill=RGBColor(0xFF, 0xF4, 0xEC))
    add_rect(slide, Inches(0.45), y2, Inches(0.15), Inches(1.95),
             fill=ACCENT)
    add_text(slide, Inches(0.75), y2 + Inches(0.08), Inches(11), Inches(0.4),
             "改进切入点 (与老师多轮讨论后确定)",
             size=16, bold=True, color=ACCENT)
    add_text(slide, Inches(0.75), y2 + Inches(0.55), Inches(11.8), Inches(1.4),
             "① 引入 文生图代理 (Visual Proxy) , 为纯文本推理补上 “图像级校验”\n"
             "② 设计 V7 反幻觉 CoT: 代理图仅作 “诊断” , 禁止写入虚构细节\n"
             "③ 在 CLIP 特征空间做 Description Ensemble, 让 D1 兜底, D2 增量修正 —— 保证不退化",
             size=14, color=DARK_TXT)
    footer(slide)


# ==================================================
# 03 方法: 三路融合总览 (使用已有 fig_pipeline.png)
# ==================================================
def slide_method_overview(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "03  方法 · 三路融合总体框架",
                "Visual Proxy + V7 Anti-Hallucination CoT + Description Ensemble",
                page=5, total=total)

    # 左图 (pipeline)
    img_path = os.path.join(FIG, "fig_pipeline.png")
    if os.path.exists(img_path):
        add_image_fit(slide, img_path,
                      cx=int(Inches(4.6)), cy=int(Inches(4.2)),
                      max_w=int(Inches(8.5)), max_h=int(Inches(5.2)))
    # 右侧: 公式
    x = Inches(9.2)
    add_rect(slide, x, Inches(1.3), Inches(3.8), Inches(5.8),
             fill=LIGHT_BG)
    add_text(slide, x + Inches(0.2), Inches(1.4), Inches(3.5), Inches(0.4),
             "检索公式", size=16, bold=True, color=TJ_PURPLE)
    add_text(slide, x + Inches(0.2), Inches(1.8), Inches(3.6), Inches(1.8),
             "f_text =\n"
             "  norm( β · CLIP(D₁) +\n"
             "       (1−β) · CLIP(D₂) )\n\n"
             "score =\n"
             "  α · sim(f_text, g)\n"
             " + (1−α) · sim(CLIP(proxy), g)",
             size=12, color=DARK_TXT, font="Consolas")

    add_text(slide, x + Inches(0.2), Inches(4.1), Inches(3.5), Inches(0.35),
             "默认参数", size=14, bold=True, color=TJ_PURPLE)
    add_text(slide, x + Inches(0.2), Inches(4.45), Inches(3.6), Inches(1.2),
             "• β = 0.7  (D₁ 权重)\n"
             "• α = 0.9  (文本信号主导)\n"
             "• CLIP: ViT-L/14\n"
             "• MLLM: Qwen-VL-Max\n"
             "• T2I: MiniMax image-01",
             size=12, color=DARK_TXT)

    add_text(slide, x + Inches(0.2), Inches(5.75), Inches(3.5), Inches(0.35),
             "整体特点", size=14, bold=True, color=TJ_PURPLE)
    add_text(slide, x + Inches(0.2), Inches(6.1), Inches(3.6), Inches(1.0),
             "• 零训练, 纯推理\n"
             "• 三信号互补\n"
             "• D₁ 兜底防退化",
             size=12, color=DARK_TXT)
    footer(slide)


# ==================================================
# 04 方法: Visual Proxy
# ==================================================
def slide_method_proxy(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "03.1  创新点 ① · Visual Proxy",
                "将文本推理 “画” 出来, 为检索提供图像空间信号",
                page=6, total=total)

    # 左: 思路描述
    add_text(slide, Inches(0.45), Inches(1.3), Inches(6.2), Inches(0.5),
             "动机", size=18, bold=True, color=TJ_PURPLE)
    add_text(slide, Inches(0.45), Inches(1.8), Inches(6.3), Inches(2.5),
             "原方法全程在文本空间操作,\n"
             "描述一旦偏差, 下游 CLIP 检索将\n直接失败且无可恢复路径.\n"
             "\n解决: 让 MLLM 的初始描述 D₁\n经 T2I 模型 (MiniMax image-01)\n再生成一张代理图 proxy,\n然后把代理图作为图库检索时的\n“第二路” 图像查询, 与文本信号互补.",
             size=14, color=DARK_TXT)

    add_text(slide, Inches(0.45), Inches(4.7), Inches(6.2), Inches(0.5),
             "两种使用方式 (消融)", size=18, bold=True, color=TJ_PURPLE)
    add_text(slide, Inches(0.45), Inches(5.2), Inches(6.3), Inches(1.8),
             "• Plan A — 后融合: 在检索得分上做加权\n"
             "• Plan B — 前融合: 把 proxy 喂回 MLLM 做二轮精炼\n"
             "   → 最终方案 两者结合 (A 做检索 · B 做精炼)",
             size=14, color=DARK_TXT)

    # 右: 小样本验证
    x = Inches(7.0)
    add_rect(slide, x, Inches(1.3), Inches(6.0), Inches(5.9),
             fill=RGBColor(0xF8, 0xEB, 0xED))
    add_text(slide, x + Inches(0.3), Inches(1.4), Inches(5.5), Inches(0.4),
             "FashionIQ dress 子集 (50 样本) R@10", size=16, bold=True,
             color=ACCENT)
    # 简易柱状 (用矩形)
    bars = [
        ("Baseline (仅 D1)", 18.0, GRAY_TXT),
        ("Plan A 后融合 α=0.8", 26.0, TJ_PURPLE),
        ("Plan B 前融合", 22.0, DEEP_BLUE),
    ]
    max_val = 30.0
    bar_w_max = Inches(3.8)
    for i, (name, v, c) in enumerate(bars):
        y = Inches(2.0) + Inches(1.15) * i
        add_text(slide, x + Inches(0.3), y, Inches(5.5), Inches(0.35),
                 name, size=13, bold=True, color=DARK_TXT)
        # 条底
        add_rect(slide, x + Inches(0.3), y + Inches(0.45),
                 bar_w_max, Inches(0.45),
                 fill=RGBColor(0xEE, 0xEE, 0xEE))
        # 条
        bar_w = int(bar_w_max * (v / max_val))
        add_rect(slide, x + Inches(0.3), y + Inches(0.45),
                 bar_w, Inches(0.45), fill=c)
        add_text(slide, x + Inches(0.3) + bar_w_max + Inches(0.1),
                 y + Inches(0.45), Inches(1.2), Inches(0.45),
                 f"{v:.1f}", size=14, bold=True, color=c,
                 anchor=MSO_ANCHOR.MIDDLE)

    add_rect(slide, x + Inches(0.3), Inches(5.9), Inches(5.5), Inches(1.2),
             fill=WHITE, line=ACCENT, line_w=Pt(1.2))
    add_text(slide, x + Inches(0.45), Inches(5.95), Inches(5.3), Inches(1.1),
             "→ 两种方式均显著提升 baseline\n"
             "→ 验证了 proxy 确实携带可用的图像级信息",
             size=13, color=DARK_TXT)
    footer(slide)


# ==================================================
# 05 方法: V7 Anti-Hallucination (prompt evolution)
# ==================================================
def slide_method_v7(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "03.2  创新点 ② · V7 反幻觉 CoT Prompt",
                "Prompt 迭代 4 轮, 终于解决 “代理图是参照又是幻觉源” 的矛盾",
                page=7, total=total)

    # 左上: prompt evolution figure
    img_path = os.path.join(FIG, "fig_prompt_evolution.png")
    if os.path.exists(img_path):
        add_image_fit(slide, img_path,
                      cx=int(Inches(4.3)), cy=int(Inches(4.3)),
                      max_w=int(Inches(8.2)), max_h=int(Inches(5.2)))

    # 右: 关键设计
    x = Inches(8.8)
    add_rect(slide, x, Inches(1.25), Inches(4.3), Inches(5.9),
             fill=LIGHT_BG)
    add_text(slide, x + Inches(0.2), Inches(1.35), Inches(4), Inches(0.4),
             "V7 的三个关键设计", size=15, bold=True, color=TJ_PURPLE)
    items = [
        ("① 代理图 AI 生成",
         "prompt 中显式声明: 代理图含幻觉细节\n禁止写入代理图中新增物体 / 背景"),
        ("② 代理图仅做 “诊断”",
         "唯一用途: 检查颜色 / 物体 / 属性\n是否已按修改文本改变"),
        ("③ 强制短描述",
         "输出长度 ≤ 修改文本长度\n避免描述膨胀导致 CLIP 漂移"),
    ]
    y = Inches(1.75)
    for title, body in items:
        add_text(slide, x + Inches(0.2), y, Inches(4), Inches(0.4),
                 title, size=13, bold=True, color=ACCENT)
        add_text(slide, x + Inches(0.25), y + Inches(0.4), Inches(4.0),
                 Inches(1.2), body, size=11, color=DARK_TXT)
        y += Inches(1.55)

    add_text(slide, x + Inches(0.2), Inches(6.55), Inches(4), Inches(0.5),
             "CIRR R@10:  67.0 → 73.5 (+6.5)",
             size=13, bold=True, color=GREEN)
    footer(slide)


# ==================================================
# 06 方法: Description Ensemble
# ==================================================
def slide_method_ensemble(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "03.3  创新点 ③ · Description Ensemble",
                "用 D₁ 兜底 · 用 D₂ 增量修正 · 保证不退化",
                page=8, total=total)

    # 左: 思路
    add_text(slide, Inches(0.45), Inches(1.25), Inches(6.2), Inches(0.5),
             "为什么需要 Ensemble ?", size=18, bold=True, color=TJ_PURPLE)
    add_text(slide, Inches(0.45), Inches(1.75), Inches(6.4), Inches(2.1),
             "V7 在 CIRR 提升 +6.5, 但在 FashionIQ shirt 反而 −4.0\n"
             "原因: 单一 prompt 策略存在固有矛盾:\n"
             "  • 激进精炼  → 修正偏差但可能丢失细节\n"
             "  • 保守精炼  → 保留细节但改善有限\n"
             "→ 根本问题: “全或无” 替换 D₁→D₂ 粒度太粗",
             size=13, color=DARK_TXT)

    add_text(slide, Inches(0.45), Inches(4.0), Inches(6.2), Inches(0.5),
             "解决方案", size=18, bold=True, color=TJ_PURPLE)
    add_rect(slide, Inches(0.45), Inches(4.5), Inches(6.4), Inches(1.2),
             fill=LIGHT_BG, line=TJ_PURPLE, line_w=Pt(1.2))
    add_text(slide, Inches(0.6), Inches(4.55), Inches(6), Inches(1.1),
             "f_text = normalize( β · CLIP(D₁) + (1−β) · CLIP(D₂) )\n"
             "        β = 0.7      (D₁ 特征 70% · D₂ 特征 30%)",
             size=14, bold=True, color=DARK_TXT, font="Consolas")

    add_text(slide, Inches(0.45), Inches(5.9), Inches(6.5), Inches(1.5),
             "→ 在 CLIP 特征空间做 weighted-sum 而非硬替换\n"
             "→ 即使 D₂ 在某些样本上差, D₁ 的信号仍占主导\n"
             "→ 实现所有子集的 统一正向提升",
             size=13, color=DARK_TXT)

    # 右: 对比表
    x = Inches(7.3)
    add_rect(slide, x, Inches(1.25), Inches(5.7), Inches(5.9),
             fill=WHITE, line=TJ_PURPLE, line_w=Pt(1.2))
    add_text(slide, x + Inches(0.2), Inches(1.35), Inches(5.4), Inches(0.4),
             "关键对比 (R@10)", size=15, bold=True, color=TJ_PURPLE)

    # 表头
    tx = [x + Inches(0.2), x + Inches(2.4),
          x + Inches(3.7), x + Inches(4.7)]
    headers = ["方法", "CIRR", "shirt", "是否全正向"]
    for i, h in enumerate(headers):
        add_rect(slide, tx[i], Inches(1.85),
                 Inches(1.2) if i != 0 else Inches(2.2),
                 Inches(0.4), fill=TJ_PURPLE)
        add_text(slide, tx[i], Inches(1.88),
                 Inches(1.2) if i != 0 else Inches(2.2),
                 Inches(0.4), h, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
    rows = [
        ("Baseline (仅 D1)", "67.0", "30.0", "—", DARK_TXT),
        ("V7 (D2 替换 D1)", "73.5 (+6.5)", "26.0 (−4.0)", "否", ACCENT),
        ("V7 + Ensemble β=0.7", "69.0 (+2.0)", "30.5 (+0.5)", "是", GREEN),
    ]
    for i, row in enumerate(rows):
        ry = Inches(2.25) + Inches(0.6) * i
        name, c1, c2, c3, color = row
        add_rect(slide, tx[0], ry, Inches(2.2), Inches(0.6),
                 fill=RGBColor(0xFA, 0xFA, 0xFA) if i % 2 else WHITE)
        add_rect(slide, tx[1], ry, Inches(1.2), Inches(0.6),
                 fill=RGBColor(0xFA, 0xFA, 0xFA) if i % 2 else WHITE)
        add_rect(slide, tx[2], ry, Inches(1.0), Inches(0.6),
                 fill=RGBColor(0xFA, 0xFA, 0xFA) if i % 2 else WHITE)
        add_rect(slide, tx[3], ry, Inches(1.0), Inches(0.6),
                 fill=RGBColor(0xFA, 0xFA, 0xFA) if i % 2 else WHITE)
        add_text(slide, tx[0], ry + Inches(0.15), Inches(2.2), Inches(0.4),
                 name, size=11, bold=True, color=DARK_TXT,
                 align=PP_ALIGN.CENTER)
        add_text(slide, tx[1], ry + Inches(0.15), Inches(1.2), Inches(0.4),
                 c1, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, tx[2], ry + Inches(0.15), Inches(1.0), Inches(0.4),
                 c2, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, tx[3], ry + Inches(0.15), Inches(1.0), Inches(0.4),
                 c3, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)

    add_rect(slide, x + Inches(0.2), Inches(4.3), Inches(5.4), Inches(1.5),
             fill=RGBColor(0xEC, 0xF8, 0xEC), line=GREEN, line_w=Pt(1.2))
    add_text(slide, x + Inches(0.35), Inches(4.38), Inches(5.2), Inches(1.4),
             "Ensemble 放弃 V7 在 CIRR 上的极端提升,\n"
             "换取所有子集的统一正向提升 ——\n"
             "这才是 “工程可用” 的改进方案",
             size=13, color=DARK_TXT)
    footer(slide)


# ==================================================
# 07 实验: 设置
# ==================================================
def slide_exp_setup(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "04.1  实验设置", "9 个标准基准 · 统一评估口径",
                page=9, total=total)

    # 左: 数据集
    add_text(slide, Inches(0.5), Inches(1.25), Inches(6.2), Inches(0.4),
             "数据集 (9 个基准)", size=18, bold=True, color=TJ_PURPLE)

    datasets = [
        ("FashionIQ", "dress 1918 · shirt 1996 · toptee 1923", "R@10 / R@50"),
        ("CIRCO", "220 queries · 123K 图库 (COCO 扩展)", "mAP@5/10/25/50"),
        ("CIRR", "4181 queries · 2297 图库", "R@1/5/10/50 · R_sub"),
        ("GeneCIS", "change/focus × object/attribute (4 子集)", "R@1/2/3"),
    ]
    for i, (n, detail, m) in enumerate(datasets):
        y = Inches(1.8) + Inches(0.95) * i
        add_rect(slide, Inches(0.5), y, Inches(6.2), Inches(0.85),
                 fill=LIGHT_BG)
        add_rect(slide, Inches(0.5), y, Inches(0.15), Inches(0.85),
                 fill=TJ_PURPLE)
        add_text(slide, Inches(0.75), y + Inches(0.05), Inches(4), Inches(0.3),
                 n, size=14, bold=True, color=DEEP_BLUE)
        add_text(slide, Inches(0.75), y + Inches(0.35), Inches(5.5), Inches(0.25),
                 detail, size=11, color=GRAY_TXT)
        add_text(slide, Inches(0.75), y + Inches(0.58), Inches(5.5), Inches(0.25),
                 "指标: " + m, size=11, bold=True, color=ACCENT)

    # 右: 技术栈
    x = Inches(7.1)
    add_text(slide, x, Inches(1.25), Inches(5.9), Inches(0.4),
             "技术栈与资源", size=18, bold=True, color=TJ_PURPLE)

    stacks = [
        ("MLLM", "Qwen-VL-Max (阿里云 DashScope)",
         "原论文 GPT-4o → 因 API 可用性替换\n差异仅影响 baseline 绝对值, 不影响相对改进"),
        ("T2I 模型", "MiniMax image-01",
         "初始描述 D₁ → 代理图 proxy · 512×512"),
        ("视觉编码器", "CLIP ViT-L/14 (与论文一致)",
         "同时编码 D₁ / D₂ / proxy / gallery"),
        ("计算资源", "Linux 3.6GB RAM  +  Windows RTX 4060 8GB",
         "MLLM / T2I 云端 API · CLIP 特征本地 GPU 计算"),
    ]
    for i, (k, v, note) in enumerate(stacks):
        y = Inches(1.8) + Inches(1.25) * i
        add_rect(slide, x, y, Inches(5.9), Inches(1.15),
                 fill=WHITE, line=RGBColor(0xCC, 0xCC, 0xCC))
        add_text(slide, x + Inches(0.15), y + Inches(0.05),
                 Inches(1.2), Inches(0.4), k, size=13, bold=True,
                 color=TJ_PURPLE)
        add_text(slide, x + Inches(1.4), y + Inches(0.05),
                 Inches(4.4), Inches(0.4), v, size=13, bold=True,
                 color=DARK_TXT)
        add_text(slide, x + Inches(0.15), y + Inches(0.5),
                 Inches(5.7), Inches(0.6), note, size=10, color=GRAY_TXT)

    footer(slide)


# ==================================================
# 08 实验: 主结果 (使用 fig_main_results.png)
# ==================================================
def slide_exp_main(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "04.2  主结果 · 9 个数据集全量实验",
                "35 / 35 指标全部提升 · 9 / 9 主指标全部正向",
                page=10, total=total)

    # 主图
    img_path = os.path.join(FIG, "fig_main_results.png")
    if os.path.exists(img_path):
        add_image_fit(slide, img_path,
                      cx=int(Inches(4.5)), cy=int(Inches(4.3)),
                      max_w=int(Inches(8.5)), max_h=int(Inches(5.3)))

    # 右侧数据
    x = Inches(9.0)
    add_rect(slide, x, Inches(1.25), Inches(4.1), Inches(5.9),
             fill=LIGHT_BG)
    add_text(slide, x + Inches(0.15), Inches(1.35), Inches(3.9),
             Inches(0.4), "关键提升 (主指标)", size=15, bold=True,
             color=TJ_PURPLE)

    rows = [
        ("FIQ dress R@10", "+3.49", "+22.1%"),
        ("FIQ toptee R@10", "+4.26", "+18.4%"),
        ("CIRCO mAP@10", "+5.05", "+31.2%"),
        ("CIRR R@1", "+2.94", "+12.8%"),
        ("GeneCIS ch_obj R@1", "+11.68", "+84.5%"),
        ("GeneCIS ch_attr R@1", "+9.14", "+72.3%"),
        ("GeneCIS fo_attr R@1", "+9.01", "+47.9%"),
    ]
    y0 = Inches(1.8)
    for i, (n, d, r) in enumerate(rows):
        y = y0 + Inches(0.48) * i
        add_text(slide, x + Inches(0.15), y, Inches(2.2), Inches(0.4),
                 n, size=11, bold=True, color=DARK_TXT)
        add_text(slide, x + Inches(2.35), y, Inches(0.8), Inches(0.4),
                 d, size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
        add_text(slide, x + Inches(3.2), y, Inches(0.85), Inches(0.4),
                 r, size=11, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)

    add_rect(slide, x + Inches(0.15), Inches(5.3), Inches(3.8), Inches(1.7),
             fill=WHITE, line=GREEN, line_w=Pt(1.5))
    add_text(slide, x + Inches(0.25), Inches(5.4), Inches(3.6), Inches(1.6),
             "✓ 35 / 35 指标全部提升\n"
             "✓ 无持平、无微降\n"
             "✓ GeneCIS 相对提升 47%~85%\n"
             "✓ 零训练 · 纯推理改进",
             size=13, bold=True, color=GREEN)
    footer(slide)


# ==================================================
# 09 实验: FashionIQ 细化
# ==================================================
def slide_exp_fiq(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "04.3  FashionIQ 细化结果",
                "服装时尚属性 · 3 子集 12 指标全部提升",
                page=11, total=total)

    img_path = os.path.join(FIG, "fig_fashioniq.png")
    if os.path.exists(img_path):
        add_image_fit(slide, img_path,
                      cx=int(Inches(4.3)), cy=int(Inches(4.3)),
                      max_w=int(Inches(8.0)), max_h=int(Inches(5.2)))

    # 右: 说明
    x = Inches(8.6)
    add_rect(slide, x, Inches(1.25), Inches(4.5), Inches(5.95),
             fill=LIGHT_BG)
    add_text(slide, x + Inches(0.2), Inches(1.35), Inches(4.2),
             Inches(0.4), "关键观察", size=15, bold=True, color=TJ_PURPLE)
    items = [
        ("dress 子集",
         "R@10 +3.49 (22.1%)\nR@50 +5.42 (16.6%)\n→ proxy 对长描述效果显著"),
        ("shirt 子集",
         "R@10 +1.35 (5.2%)\n→ V7 短描述策略原本有退化风险\n→ Ensemble 成功 “兜底”"),
        ("toptee 子集",
         "R@10 +4.26 (18.4%)\nR@50 +5.25 (12.7%)\n→ 改进幅度仅次于 dress"),
    ]
    y = Inches(1.85)
    for t, b in items:
        add_text(slide, x + Inches(0.25), y, Inches(4.2), Inches(0.35),
                 t, size=13, bold=True, color=ACCENT)
        add_text(slide, x + Inches(0.35), y + Inches(0.4), Inches(4.2),
                 Inches(1.2), b, size=11, color=DARK_TXT)
        y += Inches(1.55)

    add_text(slide, x + Inches(0.2), Inches(6.65), Inches(4.2),
             Inches(0.4), "与原论文差距主要来自 MLLM 替换, 不影响改进有效性",
             size=10, color=GRAY_TXT)
    footer(slide)


# ==================================================
# 10 消融: α/β 网格搜索 heatmap
# ==================================================
def slide_ablation_grid(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "05.1  消融 · α/β 网格搜索",
                "GeneCIS 各子集最优参数差异显著 · task-adaptive 有价值",
                page=12, total=total)

    img_path = os.path.join(FIG, "fig_heatmap.png")
    if os.path.exists(img_path):
        add_image_fit(slide, img_path,
                      cx=int(Inches(4.3)), cy=int(Inches(4.3)),
                      max_w=int(Inches(8.2)), max_h=int(Inches(5.2)))

    x = Inches(8.8)
    add_rect(slide, x, Inches(1.25), Inches(4.3), Inches(5.9),
             fill=LIGHT_BG)
    add_text(slide, x + Inches(0.2), Inches(1.35), Inches(4), Inches(0.4),
             "超参数含义", size=14, bold=True, color=TJ_PURPLE)
    add_text(slide, x + Inches(0.25), Inches(1.75), Inches(4.1),
             Inches(1.4),
             "β:  D₁ 特征权重 (1→仅 D₁, 0→仅 D₂)\n"
             "α:  文本信号 vs. 代理图信号\n     (1→仅文本, 0→仅代理图)",
             size=11, color=DARK_TXT)

    add_text(slide, x + Inches(0.2), Inches(3.15), Inches(4), Inches(0.4),
             "GeneCIS 各子集最优 (α, β)", size=14, bold=True,
             color=TJ_PURPLE)
    rows = [
        ("change_object", "α=0.95  β=0.50", "+0.77"),
        ("focus_object", "α=0.80  β=1.00", "+0.82"),
        ("change_attr", "(需重建 gallery)", "—"),
        ("focus_attr", "α=1.00  β=0.60", "+1.05"),
    ]
    y = Inches(3.6)
    for n, cfg, d in rows:
        add_text(slide, x + Inches(0.25), y, Inches(1.8), Inches(0.35),
                 n, size=11, bold=True, color=DARK_TXT)
        add_text(slide, x + Inches(1.95), y, Inches(1.7), Inches(0.35),
                 cfg, size=10, color=DEEP_BLUE, font="Consolas")
        add_text(slide, x + Inches(3.5), y, Inches(0.7), Inches(0.35),
                 d, size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
        y += Inches(0.48)

    add_rect(slide, x + Inches(0.2), Inches(5.7), Inches(4), Inches(1.4),
             fill=WHITE, line=ACCENT, line_w=Pt(1.2))
    add_text(slide, x + Inches(0.3), Inches(5.75), Inches(3.9),
             Inches(1.35),
             "结论:\n• focus_object 需 高 β + 低 α (反向)\n"
             "• 其他 3 子集偏 低 β + 高 α\n"
             "→ 证明 task-adaptive 调参的必要性",
             size=11, color=DARK_TXT)
    footer(slide)


# ==================================================
# 11 消融: 组件级
# ==================================================
def slide_ablation_components(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "05.2  消融 · 三个组件的独立贡献",
                "Visual Proxy / V7 CoT / Ensemble 逐个对比基线",
                page=13, total=total)

    img_path = os.path.join(FIG, "fig_ablation.png")
    if os.path.exists(img_path):
        add_image_fit(slide, img_path,
                      cx=int(Inches(4.5)), cy=int(Inches(4.3)),
                      max_w=int(Inches(8.6)), max_h=int(Inches(5.3)))

    x = Inches(9.3)
    add_rect(slide, x, Inches(1.25), Inches(3.85), Inches(5.9),
             fill=LIGHT_BG)
    add_text(slide, x + Inches(0.2), Inches(1.35), Inches(3.6),
             Inches(0.4), "消融分析 (CIRR R@10)", size=14, bold=True,
             color=TJ_PURPLE)
    rows = [
        ("Baseline (D1 only)", "67.0"),
        ("+ Visual Proxy", "69.5  (+2.5)"),
        ("+ V7 CoT (D1→D2)", "73.5  (+6.5)"),
        ("+ V7 + Proxy (无Ens)", "70.5  (+3.5)"),
        ("Full (三路融合)", "68.7  (+1.7)*"),
    ]
    y = Inches(1.9)
    for n, v in rows:
        add_text(slide, x + Inches(0.25), y, Inches(2.5), Inches(0.35),
                 n, size=11, bold=True, color=DARK_TXT)
        add_text(slide, x + Inches(2.8), y, Inches(0.95), Inches(0.35),
                 v, size=11, bold=True, color=GREEN,
                 align=PP_ALIGN.RIGHT)
        y += Inches(0.5)

    add_text(slide, x + Inches(0.2), Inches(4.6), Inches(3.6),
             Inches(0.35), "* 牺牲 CIRR 单点峰值, 换取 9 子集全正向",
             size=10, color=GRAY_TXT)

    add_rect(slide, x + Inches(0.2), Inches(5.1), Inches(3.5),
             Inches(2.0), fill=WHITE, line=TJ_PURPLE, line_w=Pt(1.2))
    add_text(slide, x + Inches(0.3), Inches(5.18), Inches(3.4),
             Inches(1.9),
             "关键结论:\n• V7 贡献最大 (+6.5)\n"
             "• Proxy 与 Ensemble 协同\n  使结果稳健\n"
             "• 三者缺一不可",
             size=11, color=DARK_TXT)
    footer(slide)


# ==================================================
# 12 相对提升图 (fig_relative.png)
# ==================================================
def slide_relative(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "04.4  相对提升概览", "9 个数据集主指标相对提升柱状图",
                page=14, total=total)

    img_path = os.path.join(FIG, "fig_relative.png")
    if os.path.exists(img_path):
        add_image_fit(slide, img_path,
                      cx=int(Inches(4.6)), cy=int(Inches(4.3)),
                      max_w=int(Inches(9.0)), max_h=int(Inches(5.4)))

    # 右侧 TOP3
    x = Inches(9.8)
    add_rect(slide, x, Inches(1.3), Inches(3.35), Inches(5.85),
             fill=LIGHT_BG)
    add_text(slide, x + Inches(0.15), Inches(1.4), Inches(3.2),
             Inches(0.4), "TOP 3 相对提升", size=14, bold=True,
             color=TJ_PURPLE)

    tops = [
        ("1", "GeneCIS ch_obj", "+84.5%", "13.83 → 25.51"),
        ("2", "GeneCIS ch_attr", "+72.3%", "12.65 → 21.79"),
        ("3", "GeneCIS fo_attr", "+47.9%", "18.82 → 27.83"),
    ]
    y = Inches(1.95)
    for r, n, v, delta in tops:
        add_rect(slide, x + Inches(0.15), y, Inches(3.1),
                 Inches(1.0), fill=WHITE, line=ACCENT, line_w=Pt(1.0))
        add_text(slide, x + Inches(0.2), y + Inches(0.08), Inches(0.4),
                 Inches(0.4), r, size=18, bold=True, color=ACCENT)
        add_text(slide, x + Inches(0.6), y + Inches(0.08), Inches(2.5),
                 Inches(0.4), n, size=12, bold=True, color=DEEP_BLUE)
        add_text(slide, x + Inches(0.6), y + Inches(0.45), Inches(2.5),
                 Inches(0.3), delta, size=10, color=GRAY_TXT)
        add_text(slide, x + Inches(2.35), y + Inches(0.25), Inches(0.85),
                 Inches(0.45), v, size=12, bold=True, color=GREEN,
                 align=PP_ALIGN.RIGHT)
        y += Inches(1.2)

    add_text(slide, x + Inches(0.2), Inches(5.8), Inches(3.2),
             Inches(1.4),
             "→ GeneCIS 受益最大\n"
             "→ MLLM 推理 + proxy 校验 对\n  细粒度属性/物体修改任务\n  特别有效",
             size=11, color=DARK_TXT)
    footer(slide)


# ==================================================
# 13 贡献总结
# ==================================================
def slide_contribution(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "06.1  工作贡献", "三层创新 · 1 个统一工程方案",
                page=15, total=total)

    items = [
        ("C1", "方法学",
         "提出 三路融合 ZS-CIR 框架: Visual Proxy (图像信号) + V7 Anti-Hallucination CoT (描述校验) + Description Ensemble (特征空间兜底), 首次系统性结合三种机制用于零样本 CIR",
         TJ_PURPLE),
        ("C2", "Prompt 工程",
         "通过 4 轮迭代, 诊断并解决 “代理图既是参照又是幻觉源” 的核心矛盾,\n提出 “诊断而非描述” 的设计原则, 将 MLLM 从自由生成器降维为结构化校验器",
         ACCENT),
        ("C3", "实验验证",
         "在 9 个标准基准上完成全量实验, 35/35 指标全部提升;\n"
         "FashionIQ 相对提升 5~22%, CIRCO mAP 相对提升 31.2%, GeneCIS R@1 提升 47~85%",
         DEEP_BLUE),
        ("C4", "工程与可复现",
         "完整开源代码 / LaTeX 论文 / 配图代码 / 全量实验 JSON · 已上线 GitHub\n"
         "提供 Qwen-VL + MiniMax 的低成本替代方案, 成本 ≈ GPT-4o 方案的 1/10",
         GREEN),
    ]
    y = Inches(1.35)
    for code, title, body, color in items:
        add_rect(slide, Inches(0.5), y, Inches(0.9), Inches(1.3),
                 fill=color)
        add_text(slide, Inches(0.5), y + Inches(0.15), Inches(0.9),
                 Inches(0.5), code, size=26, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(slide, Inches(0.5), y + Inches(0.7), Inches(0.9),
                 Inches(0.4), "贡献", size=10, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(1.55), y, Inches(11.35), Inches(1.3),
                 fill=LIGHT_BG)
        add_text(slide, Inches(1.75), y + Inches(0.1), Inches(11),
                 Inches(0.45), title, size=16, bold=True, color=color)
        add_text(slide, Inches(1.75), y + Inches(0.55), Inches(11),
                 Inches(0.75), body, size=12, color=DARK_TXT)
        y += Inches(1.4)
    footer(slide)


# ==================================================
# 14 局限性与未来工作
# ==================================================
def slide_future(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    page_header(slide, "06.2  局限性与未来工作",
                "工程与学术两个维度的后续拓展",
                page=16, total=total)

    # 左: 局限
    add_rect(slide, Inches(0.45), Inches(1.25), Inches(6.2), Inches(5.9),
             fill=RGBColor(0xFF, 0xF1, 0xE8))
    add_text(slide, Inches(0.65), Inches(1.35), Inches(5.9),
             Inches(0.4), "当前工作局限", size=18, bold=True,
             color=ORANGE)

    limits = [
        ("推理成本",
         "每条 query 需 2 次 MLLM 调用 + 1 次 T2I 调用,\n耗时约 15-25s, 尚未满足实时检索需求"),
        ("baseline 与原论文差距",
         "MLLM 替换为 Qwen-VL 使 FIQ 绝对值下降 13pp,\n但改进幅度仍保持一致, 方法本身不受影响"),
        ("超参数依赖",
         "α=0.9, β=0.7 在多数任务最优, 但 GeneCIS\nfocus_object 需反向, 需 task-adaptive 调参"),
        ("proxy 质量波动",
         "T2I 模型对长描述 / 复杂场景 仍会生成低质代理图,\n部分样本 proxy 反而成为干扰"),
    ]
    y = Inches(1.85)
    for t, b in limits:
        add_text(slide, Inches(0.7), y, Inches(5.8), Inches(0.4),
                 f"• {t}", size=14, bold=True, color=ORANGE)
        add_text(slide, Inches(1.0), y + Inches(0.4), Inches(5.4),
                 Inches(0.9), b, size=11, color=DARK_TXT)
        y += Inches(1.25)

    # 右: 未来
    add_rect(slide, Inches(6.9), Inches(1.25), Inches(6.05), Inches(5.9),
             fill=RGBColor(0xEC, 0xF5, 0xFF))
    add_text(slide, Inches(7.1), Inches(1.35), Inches(5.9),
             Inches(0.4), "后续可拓展方向", size=18, bold=True,
             color=DEEP_BLUE)
    futures = [
        ("F1 proxy 质量感知",
         "用 CLIP 相似度评估 proxy 质量, 动态调整 α"),
        ("F2 蒸馏到小模型",
         "将三路融合流程蒸馏成一个轻量 CLIP 投影头\n做到 10ms 级在线检索"),
        ("F3 迭代式 proxy 生成",
         "多轮 “文本→图像→再推理” 循环, 直到描述收敛"),
        ("F4 跨模态评估框架",
         "把 proxy 引入开放域 CIR benchmark 设计\n提供更严格的 ZS-CIR 评估口径"),
    ]
    y = Inches(1.85)
    for t, b in futures:
        add_text(slide, Inches(7.15), y, Inches(5.75), Inches(0.4),
                 f"• {t}", size=14, bold=True, color=DEEP_BLUE)
        add_text(slide, Inches(7.45), y + Inches(0.4), Inches(5.4),
                 Inches(0.9), b, size=11, color=DARK_TXT)
        y += Inches(1.25)
    footer(slide)


# ==================================================
# 15 感谢
# ==================================================
def slide_thanks(total):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, TJ_PURPLE)
    # 顶底装饰条
    add_rect(slide, 0, Inches(3.35), SW, Inches(0.08), fill=ACCENT)

    add_text(slide, Inches(0.5), Inches(1.5), Inches(12.3),
             Inches(1.4), "Thanks for Your Attention",
             size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.6),
             "恳请各位老师批评指正", size=22, color=RGBColor(0xE0, 0xD0, 0xF0),
             align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.5),
             "GitHub 开源:", size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5),
             "https://github.com/Haookok/osrcir-threeway-fusion",
             size=16, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.6),
             "杨浩铭 · 同济大学 · 计算机科学与技术",
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
             "2026 届本科毕业设计",
             size=14, color=RGBColor(0xDD, 0xCC, 0xEE),
             align=PP_ALIGN.CENTER)


# ==================================================
# MAIN
# ==================================================
def main():
    TOTAL = 16  # 将会用到的页数
    slide_cover()
    slide_toc(TOTAL)
    slide_background(TOTAL)
    slide_problem(TOTAL)
    slide_method_overview(TOTAL)
    slide_method_proxy(TOTAL)
    slide_method_v7(TOTAL)
    slide_method_ensemble(TOTAL)
    slide_exp_setup(TOTAL)
    slide_exp_main(TOTAL)
    slide_exp_fiq(TOTAL)
    slide_ablation_grid(TOTAL)
    slide_ablation_components(TOTAL)
    slide_relative(TOTAL)
    slide_contribution(TOTAL)
    slide_future(TOTAL)
    slide_thanks(TOTAL)  # 不计页码

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"[OK] saved -> {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
