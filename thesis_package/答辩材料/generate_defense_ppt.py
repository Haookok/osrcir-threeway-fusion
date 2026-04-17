"""
生成同济大学本科毕业答辩 PPT
基于 python-pptx 自动生成，包含主体幻灯片 + 备用幻灯片
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(OUT_DIR, '..', '图表及代码')

# ============ 同济配色 ============
TJ_BLUE = RGBColor(0x1F, 0x40, 0x90)       # 同济蓝
TJ_RED = RGBColor(0xC8, 0x28, 0x28)        # 强调红
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x80, 0x80, 0x80)
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
ACCENT_GREEN = RGBColor(0x16, 0xA3, 0x4A)
ACCENT_ORANGE = RGBColor(0xEA, 0x58, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_bar(slide, title_text, subtitle=None):
    """在每张幻灯片顶部添加统一的标题栏"""
    # 标题栏背景条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.7)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TJ_BLUE
    bar.line.fill.background()

    # 标题文字
    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = '黑体'

    # 右上角副标题（如论文题目缩写）
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(9), Inches(0.15), Inches(4), Inches(0.4)
        )
        sp = sub_box.text_frame.paragraphs[0]
        sp.alignment = PP_ALIGN.RIGHT
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.size = Pt(11)
        srun.font.color.rgb = WHITE
        srun.font.name = '宋体'


def add_page_footer(slide, page_num, total):
    """在底部添加页码"""
    # 底部分割线
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(7.0), Inches(12.8), Inches(7.0))
    line.line.color.rgb = LIGHT_GRAY
    line.line.width = Pt(0.5)

    # 左下角作者信息
    info = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(6), Inches(0.3))
    p = info.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "杨昊明  |  同济大学电子与信息工程学院  |  计算机科学与技术"
    run.font.size = Pt(10)
    run.font.color.rgb = LIGHT_GRAY
    run.font.name = '宋体'

    # 右下角页码
    pg = slide.shapes.add_textbox(Inches(12), Inches(7.1), Inches(1.2), Inches(0.3))
    pp = pg.text_frame.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    prun = pp.add_run()
    prun.text = f"{page_num} / {total}"
    prun.font.size = Pt(10)
    prun.font.color.rgb = LIGHT_GRAY
    prun.font.name = 'Times New Roman'


def add_text_box(slide, x, y, w, h, text, size=14, bold=False, color=None, align=PP_ALIGN.LEFT, font='宋体'):
    """添加文本框的便捷函数"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    if color:
        run.font.color.rgb = color
    return box, tf


def add_bullet_list(slide, x, y, w, h, items, size=16, line_spacing=1.3):
    """添加带项目符号的列表，items 可以是 str 或 (level, text) 元组"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)

    for i, item in enumerate(items):
        if isinstance(item, tuple):
            level, text = item
            if not isinstance(level, int):
                level = 1 if level == "" else 0
        else:
            level, text = 0, item

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = level
        p.line_spacing = line_spacing
        p.alignment = PP_ALIGN.LEFT

        # 添加项目符号前缀
        prefix = "●  " if level == 0 else "○  "
        run = p.add_run()
        run.text = prefix + text
        run.font.size = Pt(size - level * 2)
        run.font.color.rgb = DARK_GRAY
        run.font.name = '宋体'
    return box


def add_section_divider(slide, part_num, part_title, en_title):
    """添加章节过渡页（大号居中）"""
    # 背景色块
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = TJ_BLUE
    bg.line.fill.background()

    # 章节编号
    num_box = slide.shapes.add_textbox(Inches(0), Inches(2.2), Inches(13.333), Inches(1.2))
    p = num_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"PART  {part_num:02d}"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = 'Times New Roman'

    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(0), Inches(3.7), Inches(13.333), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = part_title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = '黑体'

    # 英文副标题
    en_box = slide.shapes.add_textbox(Inches(0), Inches(4.7), Inches(13.333), Inches(0.5))
    p = en_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = en_title
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0xCF, 0xD8, 0xE3)
    run.font.name = 'Times New Roman'

    # 底部分隔线
    line = slide.shapes.add_connector(1, Inches(5.5), Inches(5.4), Inches(7.8), Inches(5.4))
    line.line.color.rgb = WHITE
    line.line.width = Pt(2)


# ==========================================================
# 开始构建 PPT
# ==========================================================
prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 宽屏
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # 空白版式

TOTAL = 28   # 总页数（主体 + 备用），将在末尾统一设置

# =========================================================
# Slide 1: 封面
# =========================================================
s = prs.slides.add_slide(BLANK)

# 顶部蓝色条
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.5))
bar.fill.solid(); bar.fill.fore_color.rgb = TJ_BLUE; bar.line.fill.background()

# 底部蓝色条
bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5))
bar2.fill.solid(); bar2.fill.fore_color.rgb = TJ_BLUE; bar2.line.fill.background()

# 英文校名
add_text_box(s, 0, 0.9, 13.333, 0.6, "TONGJI  UNIVERSITY",
             size=26, bold=True, color=TJ_BLUE, align=PP_ALIGN.CENTER, font='Times New Roman')

# "毕业设计（论文）答辩"
add_text_box(s, 0, 1.5, 13.333, 0.6, "本科毕业设计（论文）答辩",
             size=24, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER, font='黑体')

# 分隔线
line = s.shapes.add_connector(1, Inches(4), Inches(2.3), Inches(9.333), Inches(2.3))
line.line.color.rgb = TJ_BLUE
line.line.width = Pt(2)

# 论文题目
add_text_box(s, 0.5, 2.7, 12.333, 1.0,
             "基于视觉代理与描述融合的",
             size=36, bold=True, color=TJ_BLUE, align=PP_ALIGN.CENTER, font='黑体')
add_text_box(s, 0.5, 3.7, 12.333, 1.0,
             "零样本组合式图像检索改进",
             size=36, bold=True, color=TJ_BLUE, align=PP_ALIGN.CENTER, font='黑体')

# 英文标题
add_text_box(s, 0.5, 4.9, 12.333, 0.5,
             "Improving Zero-Shot Composed Image Retrieval",
             size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, font='Times New Roman')
add_text_box(s, 0.5, 5.3, 12.333, 0.5,
             "via Visual Proxy and Description Fusion",
             size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, font='Times New Roman')

# 学生信息
add_text_box(s, 0.5, 6.05, 12.333, 0.4,
             "答辩人：杨昊明     学院：电子与信息工程学院     专业：计算机科学与技术",
             size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_text_box(s, 0.5, 6.45, 12.333, 0.4,
             "指导教师：（填写指导教师）         日期：2026 年 5 月",
             size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# =========================================================
# Slide 2: 目录
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "目  录", "Contents")

toc = [
    ("01", "研究背景与意义", "Background & Motivation"),
    ("02", "国内外研究现状与问题", "Related Work & Problem"),
    ("03", "三路融合方法", "Three-Way Fusion Method"),
    ("04", "实验结果与分析", "Experiments & Analysis"),
    ("05", "总结与展望", "Conclusion & Future Work"),
]

y0 = 1.4
for i, (num, zh, en) in enumerate(toc):
    y = y0 + i * 0.95
    # 左侧大号数字
    num_box = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.0), Inches(y), Inches(0.75), Inches(0.75))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = TJ_BLUE
    num_box.line.fill.background()
    tf = num_box.text_frame
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.size = Pt(20); r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = 'Times New Roman'

    # 中文标题
    add_text_box(s, 3.0, y + 0.05, 7, 0.45, zh, size=22, bold=True, color=DARK_GRAY, font='黑体')
    # 英文副标题
    add_text_box(s, 3.0, y + 0.48, 7, 0.35, en, size=12, color=LIGHT_GRAY, font='Times New Roman')

add_page_footer(s, 2, TOTAL)


# =========================================================
# Slide 3: PART 01 过渡页 - 研究背景
# =========================================================
s = prs.slides.add_slide(BLANK)
add_section_divider(s, 1, "研究背景与意义", "Background & Motivation")


# =========================================================
# Slide 4: 研究背景
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "研究背景  |  什么是组合式图像检索？")

# 左侧：任务定义
add_text_box(s, 0.6, 1.0, 6, 0.5, "任务定义", size=20, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 0.6, 1.5, 6.2, 2.5, [
    "输入：一张参考图像 + 一段修改文本",
    "输出：图库中满足修改意图的目标图像",
    "核心：既保留参考图关键内容，又应用文本修改",
], size=15)

add_text_box(s, 0.6, 3.8, 6, 0.5, "为什么重要？", size=20, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 0.6, 4.3, 6.2, 2.5, [
    "电商搜索：\"与这件同款但改成红色\"",
    "图像编辑检索：\"保留场景但换成狗\"",
    "传统检索无法同时理解视觉 + 语言约束",
], size=15)

# 右侧：示意图/说明
# 举例方框
ex_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(1.1), Inches(5.6), Inches(5.8))
ex_bg.fill.solid(); ex_bg.fill.fore_color.rgb = BG_LIGHT
ex_bg.line.color.rgb = TJ_BLUE; ex_bg.line.width = Pt(1.5)

add_text_box(s, 7.5, 1.25, 5.2, 0.5, "▎典型场景举例", size=18, bold=True, color=TJ_BLUE, font='黑体')

add_text_box(s, 7.5, 1.9, 5.2, 0.4, "参考图（Reference）", size=13, bold=True, color=DARK_GRAY)
ref = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(2.3), Inches(5.2), Inches(0.9))
ref.fill.solid(); ref.fill.fore_color.rgb = WHITE
ref.line.color.rgb = LIGHT_GRAY
add_text_box(s, 7.5, 2.5, 5.2, 0.5, "[ 一件蓝色连衣裙 ]", size=15, align=PP_ALIGN.CENTER, color=LIGHT_GRAY)

add_text_box(s, 7.5, 3.3, 5.2, 0.4, "修改文本（Modification）", size=13, bold=True, color=DARK_GRAY)
mod = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(3.7), Inches(5.2), Inches(0.7))
mod.fill.solid(); mod.fill.fore_color.rgb = RGBColor(0xFF, 0xF7, 0xE6)
mod.line.color.rgb = ACCENT_ORANGE
add_text_box(s, 7.5, 3.85, 5.2, 0.4, "\"改成红色，袖子变短\"", size=15, align=PP_ALIGN.CENTER, color=ACCENT_ORANGE, bold=True)

# 箭头
arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.9), Inches(4.5), Inches(0.4), Inches(0.4))
arr.fill.solid(); arr.fill.fore_color.rgb = TJ_BLUE
arr.line.fill.background()

add_text_box(s, 7.5, 5.0, 5.2, 0.4, "目标图像（Target）", size=13, bold=True, color=DARK_GRAY)
tgt = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(5.4), Inches(5.2), Inches(1.0))
tgt.fill.solid(); tgt.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE8)
tgt.line.color.rgb = ACCENT_GREEN
add_text_box(s, 7.5, 5.65, 5.2, 0.5, "[ 一件红色短袖连衣裙 ]", size=15, align=PP_ALIGN.CENTER, color=ACCENT_GREEN, bold=True)

add_page_footer(s, 4, TOTAL)


# =========================================================
# Slide 5: 研究意义（零样本设定）
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "研究意义  |  零样本组合式图像检索（ZS-CIR）")

# 对比方框：监督 vs 零样本
# 左框：监督
left_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(6.0), Inches(2.5))
left_bg.fill.solid(); left_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
left_bg.line.color.rgb = TJ_RED; left_bg.line.width = Pt(1.5)

add_text_box(s, 0.7, 1.35, 5.6, 0.5, "监督式 CIR", size=20, bold=True, color=TJ_RED, font='黑体')
add_bullet_list(s, 0.7, 1.9, 5.6, 1.8, [
    "需要目标数据集上的 <三元组> 标注",
    "训练成本高，标注稀缺",
    "无法跨域泛化到新任务",
], size=13)

# 右框：零样本
right_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(6.0), Inches(2.5))
right_bg.fill.solid(); right_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xFF)
right_bg.line.color.rgb = TJ_BLUE; right_bg.line.width = Pt(1.5)

add_text_box(s, 7.0, 1.35, 5.6, 0.5, "零样本 ZS-CIR（本文关注）", size=20, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 7.0, 1.9, 5.6, 1.8, [
    "不依赖目标数据集训练",
    "借助预训练视觉语言模型完成推理",
    "更贴合真实跨域场景需求",
], size=13)

# 下方：本课题定位
add_text_box(s, 0.6, 4.0, 12, 0.5, "▎本课题定位", size=20, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 0.6, 4.5, 12, 2.5, [
    "复现并改进 CVPR 2025 Highlight 论文 OSrCIR（基于推理的代表方法）",
    ("", "OSrCIR 用多模态大语言模型 (MLLM) 推理目标描述，再用 CLIP 编码检索"),
    "改进目标：为纯文本路径引入视觉校验 + 降低单点误差影响",
    "应用价值：电商搜索、图像编辑、内容创作辅助",
], size=14)

add_page_footer(s, 5, TOTAL)


# =========================================================
# Slide 6: PART 02 过渡页
# =========================================================
s = prs.slides.add_slide(BLANK)
add_section_divider(s, 2, "研究现状与存在的问题", "Related Work & Problem")


# =========================================================
# Slide 7: 国内外研究现状
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "国内外研究现状  |  两类主流路线")

# 路线 1
add_text_box(s, 0.5, 1.1, 12.3, 0.5, "▎路线一：基于映射的方法（Mapping-based）",
             size=18, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 0.6, 1.6, 12.2, 1.3, [
    "思路：把参考图像映射到文本空间，再与修改文本组合",
    "代表：Pic2Word (CVPR'23), SEARLE (ICCV'23), LinCIR (CVPR'24)",
    "优势：推理快  |  局限：视觉细节丢失，跨域效果差",
], size=13)

# 路线 2
add_text_box(s, 0.5, 3.25, 12.3, 0.5, "▎路线二：基于推理的方法（Reasoning-based，本文基线）",
             size=18, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 0.6, 3.75, 12.2, 2.0, [
    "思路：利用多模态大语言模型 (MLLM) 直接理解图文并生成目标描述",
    "代表：CIReVL (ICLR'24)、OSrCIR (CVPR'25 Highlight，本文基线)",
    "OSrCIR 创新：单阶段反思式思维链 (Reflective CoT)，避免两阶段信息损失",
    "仍存在的问题 → 下一页详述",
], size=13)

# 底部：本文基线
bot = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.0))
bot.fill.solid(); bot.fill.fore_color.rgb = BG_LIGHT
bot.line.color.rgb = TJ_BLUE
add_text_box(s, 0.7, 6.1, 12.0, 0.4,
             "本文以 OSrCIR 为基线展开改进",
             size=16, bold=True, color=TJ_BLUE)
add_text_box(s, 0.7, 6.45, 12.0, 0.4,
             "基线流程：参考图 + 修改文本 →（MLLM CoT 推理）→ 目标描述 → CLIP 编码 → 检索",
             size=12, color=DARK_GRAY)

add_page_footer(s, 7, TOTAL)


# =========================================================
# Slide 8: 现有方法存在的问题
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "现有方法存在的问题  |  三大瓶颈")

problems = [
    ("01", "初始描述缺乏视觉验证",
     "MLLM 仅凭参考图 + 文本一次推理生成描述，没有图像空间的反馈机制。\n" +
     "若描述对颜色/物体/属性理解错误，后续 CLIP 检索会沿着错误方向走到底。",
     TJ_RED),
    ("02", "单一路径检索较脆弱",
     "CLIP 对文本表述方式敏感。同一语义的两种表达在特征空间中位置可能相差很大。\n" +
     "只依赖一条文本路径时，整体性能易受单点误差影响。",
     ACCENT_ORANGE),
    ("03", "AI 生成图像带来的幻觉陷阱",
     "想直接让 MLLM 根据代理图重写描述：结果代理图中的虚构背景/材质被写进去，\n" +
     "描述膨胀 1.9 倍，反而远离真实目标（CIRR 性能跌 2.5~17.5 个点）。",
     ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(problems):
    y = 1.1 + i * 1.85
    # 编号圆圈
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(y + 0.15), Inches(0.8), Inches(0.8))
    circ.fill.solid(); circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    p = circ.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
    r.font.name = 'Times New Roman'

    # 标题
    add_text_box(s, 1.7, y, 11, 0.55, title, size=20, bold=True, color=color, font='黑体')
    # 描述
    add_text_box(s, 1.7, y + 0.55, 11, 1.3, desc, size=13, color=DARK_GRAY)

add_page_footer(s, 8, TOTAL)


# =========================================================
# Slide 9: 研究问题 & 本文工作
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "研究问题与本文工作")

# 上方问题框
q_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.3))
q_bg.fill.solid(); q_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF7, 0xE0)
q_bg.line.color.rgb = ACCENT_ORANGE; q_bg.line.width = Pt(2)

add_text_box(s, 0.8, 1.25, 12, 0.5, "▎关键研究问题", size=18, bold=True, color=ACCENT_ORANGE, font='黑体')
add_text_box(s, 0.8, 1.75, 12, 0.6,
             "如何将 AI 生成图像作为 \"可用但不可信\" 的辅助信号纳入检索流程，同时避免其成为新的误差来源？",
             size=16, color=DARK_GRAY)

# 下方本文工作
add_text_box(s, 0.5, 2.7, 12, 0.5, "▎本文四项主要工作", size=20, bold=True, color=TJ_BLUE, font='黑体')

works = [
    ("①", "视觉代理机制", "将第一轮描述输入文生图生成代理图，为检索提供图像空间信号"),
    ("②", "反幻觉精炼策略", "V7 Prompt：代理图只做诊断工具，禁止复制其视觉细节"),
    ("③", "描述融合与三路融合", "CLIP 空间加权融合 D₁/D₂/代理图，统一 α/β 参数化"),
    ("④", "GeneCIS 任务自适应", "专用 prompt + task-adaptive α/β，解决短文本小图库场景"),
]

for i, (no, t, d) in enumerate(works):
    y = 3.3 + i * 0.9
    # 编号
    add_text_box(s, 0.8, y, 0.6, 0.5, no, size=28, bold=True, color=TJ_BLUE, font='Times New Roman')
    # 标题
    add_text_box(s, 1.4, y + 0.05, 3.2, 0.45, t, size=17, bold=True, color=DARK_GRAY, font='黑体')
    # 描述
    add_text_box(s, 4.6, y + 0.1, 8.2, 0.45, d, size=13, color=DARK_GRAY)

add_page_footer(s, 9, TOTAL)


# =========================================================
# Slide 10: PART 03 过渡页 - 方法
# =========================================================
s = prs.slides.add_slide(BLANK)
add_section_divider(s, 3, "三路融合方法", "Three-Way Fusion")


# =========================================================
# Slide 11: 方法总体框架（用 pipeline 图）
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "方法总体框架  |  Three-Way Fusion Pipeline")

# 插入 pipeline 图
fig_path = os.path.join(FIG_DIR, 'fig_pipeline.png')
if os.path.exists(fig_path):
    s.shapes.add_picture(fig_path, Inches(0.6), Inches(0.95), width=Inches(9.5))

# 右侧文字说明
add_text_box(s, 10.3, 1.1, 2.9, 0.4, "三阶段流程", size=16, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 10.3, 1.5, 3.0, 3.5, [
    "Stage 1：MLLM 生成 D₁ → 文生图生成代理图 P",
    "Stage 2：MLLM V7 + 原图 + 代理图 → 精炼 D₂",
    "Retrieval：CLIP 编码三路 → 加权融合 → 排序",
], size=10, line_spacing=1.2)

# 关键公式小框
fm = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.3), Inches(4.7), Inches(2.9), Inches(2.2))
fm.fill.solid(); fm.fill.fore_color.rgb = BG_LIGHT
fm.line.color.rgb = TJ_BLUE

add_text_box(s, 10.4, 4.78, 2.75, 0.35, "核心公式", size=12, bold=True, color=TJ_BLUE)
add_text_box(s, 10.4, 5.15, 2.75, 0.4,
             "f_text = β·CLIP(D₁)",
             size=10, color=DARK_GRAY, font='Consolas')
add_text_box(s, 10.4, 5.45, 2.75, 0.4,
             "       + (1-β)·CLIP(D₂)",
             size=10, color=DARK_GRAY, font='Consolas')
add_text_box(s, 10.4, 5.85, 2.75, 0.4,
             "score = α·sim(f_text, g)",
             size=10, color=DARK_GRAY, font='Consolas')
add_text_box(s, 10.4, 6.15, 2.75, 0.4,
             "      + (1-α)·sim(P, g)",
             size=10, color=DARK_GRAY, font='Consolas')
add_text_box(s, 10.4, 6.5, 2.75, 0.3, "默认 β=0.7, α=0.9",
             size=10, bold=True, color=TJ_RED)

add_page_footer(s, 11, TOTAL)


# =========================================================
# Slide 12: 创新点 1 - 视觉代理
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "创新点 ①  |  视觉代理机制（Visual Proxy）")

# 左：动机
add_text_box(s, 0.5, 1.1, 6, 0.5, "▎设计动机", size=18, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 0.5, 1.6, 6, 2.5, [
    "D₁ 既要理解意图又要决定检索方向",
    "出错后整个流程无视觉验证机制",
    "→ 把 D₁ \"画\" 出来投射到图像空间",
], size=14)

# 左：作用
add_text_box(s, 0.5, 4.0, 6, 0.5, "▎两类作用", size=18, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 0.5, 4.5, 6, 2.3, [
    "检索阶段：图像特征 ↔ 文本特征 互补",
    "精炼阶段：为 MLLM 提供视觉对照",
    ("", "可验证颜色是否改变、物体是否替换、主体是否保留"),
], size=13)

# 右：小规模验证表格
add_text_box(s, 7.0, 1.1, 6, 0.5, "▎FashionIQ dress 50 样本验证", size=18, bold=True, color=TJ_BLUE, font='黑体')

# 表格
from pptx.util import Inches, Pt
tbl = s.shapes.add_table(4, 2, Inches(7.0), Inches(1.7), Inches(6.0), Inches(2.5)).table
hdr = [("方法", "R@10"),
       ("Baseline（纯 D₁）", "18.0"),
       ("Plan A 后融合 (α=0.8)", "26.0"),
       ("Plan B 前融合", "22.0")]
for i, (c1, c2) in enumerate(hdr):
    tbl.cell(i, 0).text = c1
    tbl.cell(i, 1).text = c2
    for j in range(2):
        cell = tbl.cell(i, j)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)
                r.font.name = '宋体'
                if i == 0:
                    r.font.bold = True
                    r.font.color.rgb = WHITE
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TJ_BLUE

# 右下：结论
add_text_box(s, 7.0, 4.4, 6, 0.5, "▎关键结论", size=18, bold=True, color=ACCENT_GREEN, font='黑体')
add_bullet_list(s, 7.0, 4.9, 6, 2.0, [
    "代理图虽然含噪声，但提供了有价值的检索信号",
    "后融合 (Plan A) 效果最好，+8.0",
    "最终方案融合 Plan A + Plan B（前后融合结合）",
], size=13)

add_page_footer(s, 12, TOTAL)


# =========================================================
# Slide 13: 创新点 2 - 反幻觉 prompt + 演进过程
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "创新点 ②  |  反幻觉描述精炼（V7 Prompt）")

# 左：prompt 演进图
fig_path = os.path.join(FIG_DIR, 'fig_prompt_evolution.png')
if os.path.exists(fig_path):
    s.shapes.add_picture(fig_path, Inches(0.3), Inches(1.0), width=Inches(7.0))

# 右：V7 核心原则
add_text_box(s, 7.5, 1.1, 5.5, 0.5, "▎V7 核心原则", size=18, bold=True, color=TJ_BLUE, font='黑体')

box_v7 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.6), Inches(5.5), Inches(3.0))
box_v7.fill.solid(); box_v7.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE8)
box_v7.line.color.rgb = ACCENT_GREEN

add_text_box(s, 7.7, 1.7, 5.2, 0.5,
             "「代理图只做诊断，不做描述来源」",
             size=16, bold=True, color=ACCENT_GREEN)
add_bullet_list(s, 7.7, 2.2, 5.2, 2.3, [
    "明确告知：代理图是 AI 生成，含幻觉",
    "只允许检查修改是否应用（颜色/物体/属性）",
    "禁止复制代理图中的背景、材质、环境细节",
    "强制输出简短描述（不超过修改文本长度）",
], size=11, line_spacing=1.15)

# 右下：关键发现
add_text_box(s, 7.5, 4.9, 5.5, 0.5, "▎迭代教训", size=18, bold=True, color=TJ_RED, font='黑体')
add_bullet_list(s, 7.5, 5.4, 5.5, 1.8, [
    "V6 灾难性退化：-17.5 → 信息过多反而有害",
    "V7 突破：+6.5 → 限制信息来源比增加信息更重要",
    "但 shirt 子集 -4.0 → 需要 Ensemble 救场",
], size=11, line_spacing=1.2)

add_page_footer(s, 13, TOTAL)


# =========================================================
# Slide 14: 创新点 3 - 描述融合 + 三路融合
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "创新点 ③  |  描述融合与三路融合")

# 上：动机
add_text_box(s, 0.5, 1.0, 12.3, 0.5, "▎为什么要 \"融合\" 而不是 \"替换\"", size=18, bold=True, color=TJ_BLUE, font='黑体')
add_text_box(s, 0.5, 1.5, 12.3, 1.0,
             "V7 在某些样本上会过度压缩信息。若直接用 D₂ 替换 D₁，系统性能可能因少量异常样本而整体下降。\n" +
             "解决思路：在 CLIP 特征空间加权平均——即使 D₂ 退化，D₁ 仍占主导，不会退化。",
             size=13, color=DARK_GRAY)

# 中：公式
add_text_box(s, 0.5, 2.7, 12.3, 0.5, "▎两阶段融合公式", size=18, bold=True, color=TJ_BLUE, font='黑体')

# 公式 1
f1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.2), Inches(6.0), Inches(1.6))
f1.fill.solid(); f1.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
f1.line.color.rgb = TJ_BLUE
add_text_box(s, 0.7, 3.3, 5.8, 0.4, "① 描述融合（文本空间）", size=14, bold=True, color=TJ_BLUE)
add_text_box(s, 0.7, 3.75, 5.8, 0.5,
             "f_text = normalize( β·CLIP(D₁) + (1-β)·CLIP(D₂) )",
             size=14, color=DARK_GRAY, font='Consolas')
add_text_box(s, 0.7, 4.3, 5.8, 0.4,
             "β = 0.7（D₁ 占 70%，精炼描述占 30%）",
             size=12, color=TJ_RED, bold=True)

# 公式 2
f2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.2), Inches(6.0), Inches(1.6))
f2.fill.solid(); f2.fill.fore_color.rgb = RGBColor(0xFF, 0xE8, 0xE8)
f2.line.color.rgb = TJ_RED
add_text_box(s, 6.9, 3.3, 5.8, 0.4, "② 三路融合（检索得分）", size=14, bold=True, color=TJ_RED)
add_text_box(s, 6.9, 3.75, 5.8, 0.5,
             "score = α·sim(f_text, g) + (1-α)·sim(CLIP(P), g)",
             size=13, color=DARK_GRAY, font='Consolas')
add_text_box(s, 6.9, 4.3, 5.8, 0.4,
             "α = 0.9（文本 90%，代理图 10% 作为辅助）",
             size=12, color=TJ_BLUE, bold=True)

# 下：设计哲学
add_text_box(s, 0.5, 5.1, 12.3, 0.5, "▎设计哲学", size=18, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 0.5, 5.6, 12.3, 1.4, [
    "不追求极端提升，追求稳定增益 —— 牺牲 V7 的 +6.5 换来全部 35/35 指标正向",
    "β 偏保守（0.7）：以 D₁ 兜底，避免精炼失误",
    "α 偏保守（0.9）：代理图含幻觉，只做辅助信号",
], size=13)

add_page_footer(s, 14, TOTAL)


# =========================================================
# Slide 15: PART 04 过渡页 - 实验
# =========================================================
s = prs.slides.add_slide(BLANK)
add_section_divider(s, 4, "实验结果与分析", "Experiments & Analysis")


# =========================================================
# Slide 16: 实验设置
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "实验设置")

# 数据集
add_text_box(s, 0.5, 1.0, 12, 0.5, "▎数据集（4 个基准 × 9 个子任务）", size=18, bold=True, color=TJ_BLUE, font='黑体')

datasets = [
    ("FashionIQ", "服饰检索  |  dress / shirt / toptee 三子集  |  gallery 3.6~6.2K"),
    ("CIRCO", "开放域检索  |  220 query  |  gallery 123K"),
    ("CIRR", "自然场景检索  |  4181 query  |  gallery 2297"),
    ("GeneCIS", "条件相似  |  change/focus × object/attribute 四子集  |  局部图库 ~14"),
]
for i, (name, desc) in enumerate(datasets):
    y = 1.5 + i * 0.45
    add_text_box(s, 0.8, y, 1.8, 0.4, name, size=14, bold=True, color=TJ_BLUE, font='黑体')
    add_text_box(s, 2.6, y, 10, 0.4, desc, size=12, color=DARK_GRAY)

# 实现细节
add_text_box(s, 0.5, 3.5, 12, 0.5, "▎实现细节", size=18, bold=True, color=TJ_BLUE, font='黑体')

details = [
    ("MLLM", "Qwen-VL-Max（替代原论文 GPT-4o，接口可用性原因）"),
    ("T2I 模型", "MiniMax image-01"),
    ("编码器", "CLIP ViT-L/14（与基线一致）"),
    ("默认参数", "β = 0.7，α = 0.9（GeneCIS 采用 task-adaptive）"),
    ("运行环境", "Linux 服务器（描述生成）+ Windows RTX 4060（CLIP 编码）"),
    ("评估", "FIQ/CIRCO/CIRR 用默认参数；GeneCIS 使用专用 prompt + best α/β"),
]
for i, (k, v) in enumerate(details):
    y = 4.0 + i * 0.42
    add_text_box(s, 0.8, y, 2.2, 0.4, k, size=13, bold=True, color=DARK_GRAY, font='黑体')
    add_text_box(s, 3.0, y, 10, 0.4, v, size=12, color=DARK_GRAY)

# 底部提示
note = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4))
note.fill.solid(); note.fill.fore_color.rgb = RGBColor(0xFF, 0xF7, 0xE6)
note.line.color.rgb = ACCENT_ORANGE
add_text_box(s, 0.7, 6.8, 12, 0.35,
             "⚠ 公平性说明：所有改进与基线均基于同一 Qwen-VL-Max，对比公平；仅绝对数值低于原论文 GPT-4o 结果。",
             size=11, color=DARK_GRAY)

add_page_footer(s, 16, TOTAL)


# =========================================================
# Slide 17: 主要结果 —— 9 个数据集总览
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "主要结果  |  9 个数据集全量实验")

# 左：主结果图
fig_path = os.path.join(FIG_DIR, 'fig_main_results.png')
if os.path.exists(fig_path):
    s.shapes.add_picture(fig_path, Inches(0.3), Inches(0.95), width=Inches(8.0))

# 右：核心结论方框
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.0), Inches(4.6), Inches(5.9))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE8)
box.line.color.rgb = ACCENT_GREEN; box.line.width = Pt(2)

add_text_box(s, 8.7, 1.15, 4.3, 0.5, "▎核心结论", size=18, bold=True, color=ACCENT_GREEN, font='黑体')

add_text_box(s, 8.7, 1.65, 4.3, 0.5, "35 / 35", size=44, bold=True, color=ACCENT_GREEN, font='Times New Roman', align=PP_ALIGN.CENTER)
add_text_box(s, 8.7, 2.55, 4.3, 0.4, "评估指标全部提升", size=14, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_text_box(s, 8.7, 2.95, 4.3, 0.4, "无持平  |  无退化", size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 分隔
line2 = s.shapes.add_connector(1, Inches(9.0), Inches(3.45), Inches(12.6), Inches(3.45))
line2.line.color.rgb = ACCENT_GREEN

add_text_box(s, 8.7, 3.55, 4.3, 0.4, "亮眼数据", size=14, bold=True, color=ACCENT_GREEN)
add_bullet_list(s, 8.7, 3.95, 4.3, 2.8, [
    "FIQ dress R@10: 15.80 → 19.29 (+22.1%)",
    "CIRCO mAP@10: 16.21 → 21.26 (+31.2%)",
    "CIRR R@1: 22.96 → 25.90 (+12.8%)",
    "GeneCIS ch_obj R@1: 13.83 → 25.51 (+84.5%)",
], size=11, line_spacing=1.25)

add_page_footer(s, 17, TOTAL)


# =========================================================
# Slide 18: 相对提升 + FashionIQ 详细
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "分数据集分析  |  FashionIQ 与相对提升")

# 左：FashionIQ 分子集
fig_path = os.path.join(FIG_DIR, 'fig_fashioniq.png')
if os.path.exists(fig_path):
    s.shapes.add_picture(fig_path, Inches(0.3), Inches(0.95), width=Inches(7.5))
add_text_box(s, 0.3, 4.4, 7.5, 0.4, "图 1: FashionIQ 三子集 R@1/5/10/50", size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 右：相对提升图
fig_path = os.path.join(FIG_DIR, 'fig_relative.png')
if os.path.exists(fig_path):
    s.shapes.add_picture(fig_path, Inches(8.0), Inches(0.95), width=Inches(5.3))
add_text_box(s, 8.0, 4.4, 5.3, 0.4, "图 2: 9 数据集相对提升百分比", size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 下方文字总结
add_text_box(s, 0.5, 5.0, 12.3, 0.5, "▎观察与分析", size=18, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 0.5, 5.5, 12.3, 1.5, [
    "服饰类 (FIQ)：dress/toptee 提升显著 (+18~22%)；shirt 偏保守 (+5.2%) —— 服饰风格变化语义丰富",
    "开放域 (CIRCO)：+31.2% 最大相对提升 —— 代理图对复杂语义修改价值最高",
    "GeneCIS：+47~85% 相对提升 —— task-adaptive 参数 + 专用 prompt 效果显著",
], size=13)

add_page_footer(s, 18, TOTAL)


# =========================================================
# Slide 19: 消融实验 —— prompt + 参数
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "消融实验  |  Prompt 迭代 与 α/β 参数")

# 左：prompt 迭代柱状图
fig_path = os.path.join(FIG_DIR, 'fig_prompt_evolution.png')
if os.path.exists(fig_path):
    s.shapes.add_picture(fig_path, Inches(0.3), Inches(0.95), width=Inches(6.2))
add_text_box(s, 0.3, 4.2, 6.2, 0.4, "图 3: Prompt 迭代对 CIRR R@10 的影响", size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 右：α/β 热力图
fig_path = os.path.join(FIG_DIR, 'fig_heatmap.png')
if os.path.exists(fig_path):
    s.shapes.add_picture(fig_path, Inches(6.7), Inches(0.95), width=Inches(6.5))
add_text_box(s, 6.7, 4.2, 6.5, 0.4, "图 4: GeneCIS α/β 网格搜索（★ 为最优）", size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 下方：关键发现
add_text_box(s, 0.5, 4.7, 12.3, 0.5, "▎消融关键发现", size=18, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 0.5, 5.2, 12.3, 1.8, [
    "Prompt：V6 崩溃（-17.5）vs V7 突破（+6.5）证明 \"切断幻觉路径\" 比 \"增加信息量\" 更重要",
    "各 GeneCIS 子集偏好差异极大：change_obj (β=0.3 重 D₂) vs focus_obj (β=1.0 弃 D₂)",
    "→ 说明 \"不存在通用最优参数\"，task-adaptive 是必要的",
], size=13)

add_page_footer(s, 19, TOTAL)


# =========================================================
# Slide 20: PART 05 过渡页 - 总结
# =========================================================
s = prs.slides.add_slide(BLANK)
add_section_divider(s, 5, "总结与展望", "Conclusion & Future Work")


# =========================================================
# Slide 21: 工作总结
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "工作总结")

# 四项贡献
contribs = [
    ("1", "提出视觉代理机制",
     "首次将 AI 生成代理图同时用于检索信号 + 精炼参照",
     TJ_BLUE),
    ("2", "设计反幻觉 Prompt",
     "V7 原则解决代理图幻觉引入问题，CIRR +6.5",
     ACCENT_GREEN),
    ("3", "提出描述融合 + 三路融合",
     "统一 α/β 参数化，保证 35/35 指标全部正向",
     ACCENT_ORANGE),
    ("4", "GeneCIS 任务自适应",
     "专用 prompt + task-adaptive 参数，相对提升 +47~85%",
     TJ_RED),
]

for i, (n, t, d, c) in enumerate(contribs):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.1 + row * 2.6

    # 背景
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6.1), Inches(2.3))
    bg.fill.solid(); bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.color.rgb = c; bg.line.width = Pt(2)

    # 大数字
    add_text_box(s, x + 0.3, y + 0.15, 1.2, 1.2, n, size=60, bold=True,
                 color=c, font='Times New Roman', align=PP_ALIGN.CENTER)
    # 标题
    add_text_box(s, x + 1.7, y + 0.25, 4.2, 0.6, t, size=17, bold=True, color=DARK_GRAY, font='黑体')
    # 描述
    add_text_box(s, x + 1.7, y + 0.95, 4.3, 1.3, d, size=13, color=DARK_GRAY)

# 底部总结
bot = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.8))
bot.fill.solid(); bot.fill.fore_color.rgb = TJ_BLUE
bot.line.fill.background()
add_text_box(s, 0.5, 6.45, 12.3, 0.6,
             "9 个标准任务 × 35 项评估指标 全部正向提升 —— 验证视觉代理 + 描述融合的有效性",
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='黑体')

add_page_footer(s, 21, TOTAL)


# =========================================================
# Slide 22: 未来工作
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "未来工作展望")

futures = [
    ("查询自适应参数预测",
     "训练一个轻量预测网络，根据 query 特征自动选择 α/β，减少人工调参",
     "短期 (1-2 个月)"),
    ("更强 MLLM 的可扩展性验证",
     "升级到 GPT-4o / Gemini / Qwen3-VL，观察上限提升",
     "中期"),
    ("多代理图 & 质量过滤",
     "生成多张代理图并用 CLIP 做质量评估，挑选最可靠的作为辅助信号",
     "中期"),
    ("扩展到有监督 CIR",
     "将三路融合框架迁移到监督学习场景，与训练方法结合",
     "长期"),
]

for i, (t, d, time) in enumerate(futures):
    y = 1.1 + i * 1.35
    # 左侧箭头
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.5), Inches(y + 0.3), Inches(0.7), Inches(0.7))
    ar.fill.solid(); ar.fill.fore_color.rgb = TJ_BLUE
    ar.line.fill.background()

    # 标题
    add_text_box(s, 1.4, y + 0.1, 9, 0.5, t, size=18, bold=True, color=DARK_GRAY, font='黑体')
    # 描述
    add_text_box(s, 1.4, y + 0.65, 9, 0.6, d, size=13, color=DARK_GRAY)
    # 时间框
    tbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.8), Inches(y + 0.35), Inches(2.0), Inches(0.55))
    tbox.fill.solid(); tbox.fill.fore_color.rgb = BG_LIGHT
    tbox.line.color.rgb = TJ_BLUE
    add_text_box(s, 10.8, y + 0.45, 2.0, 0.4, time, size=12, color=TJ_BLUE, align=PP_ALIGN.CENTER, bold=True)

add_page_footer(s, 22, TOTAL)


# =========================================================
# Slide 23: 谢幕页
# =========================================================
s = prs.slides.add_slide(BLANK)

# 背景
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = TJ_BLUE
bg.line.fill.background()

# 装饰圆
deco1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(5), Inches(5))
deco1.fill.solid(); deco1.fill.fore_color.rgb = RGBColor(0x2C, 0x55, 0xA3)
deco1.line.fill.background()
deco2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(5), Inches(6), Inches(6))
deco2.fill.solid(); deco2.fill.fore_color.rgb = RGBColor(0x2C, 0x55, 0xA3)
deco2.line.fill.background()

# 大字
add_text_box(s, 0, 2.5, 13.333, 1.5, "Thank  You", size=90, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font='Times New Roman')
add_text_box(s, 0, 4.1, 13.333, 0.8, "恳请各位老师批评指正", size=28, color=WHITE,
             align=PP_ALIGN.CENTER, font='黑体')
add_text_box(s, 0, 5.0, 13.333, 0.5, "Q & A", size=20, bold=True,
             color=RGBColor(0xCF, 0xD8, 0xE3), align=PP_ALIGN.CENTER, font='Times New Roman')

add_text_box(s, 0, 6.5, 13.333, 0.4,
             "杨昊明  |  同济大学 电子与信息工程学院  |  计算机科学与技术  |  2026.05",
             size=13, color=RGBColor(0xCF, 0xD8, 0xE3), align=PP_ALIGN.CENTER)


# =========================================================
# Slide 24: 备用页 1 - 基线复现公平性
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "备用 ①  |  基线复现数值与公平性说明", "Backup")

add_text_box(s, 0.5, 1.0, 12.3, 0.5, "▎为什么本文 Baseline 低于原论文？", size=18, bold=True, color=TJ_BLUE, font='黑体')
add_text_box(s, 0.5, 1.5, 12.3, 0.5,
             "因 GPT-4o API 不可用，本工作改用 Qwen-VL-Max，MLLM 能力差异导致 baseline 数值偏低。",
             size=13, color=DARK_GRAY)

# 对比表格
tbl = s.shapes.add_table(6, 5,
                          Inches(0.5), Inches(2.2),
                          Inches(12.3), Inches(3.0)).table
hdr = ["数据集", "主指标", "原论文 (GPT-4o)", "本工作 Baseline (Qwen-VL)", "备注"]
rows = [
    ["FIQ dress", "R@10", "29.70", "15.80", "直接可比"],
    ["FIQ shirt", "R@10", "33.17", "26.00", "直接可比"],
    ["FIQ toptee", "R@10", "36.92", "23.09", "直接可比"],
    ["CIRR", "R@1", "29.45", "22.96", "不同 split"],
    ["CIRCO", "mAP@10", "25.33", "16.21", "不同 split"],
]
for j, h in enumerate(hdr):
    tbl.cell(0, j).text = h
for i, r in enumerate(rows):
    for j, v in enumerate(r):
        tbl.cell(i+1, j).text = v
for i in range(6):
    for j in range(5):
        for p in tbl.cell(i, j).text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(12)
                r.font.name = '宋体'
                if i == 0:
                    r.font.bold = True; r.font.color.rgb = WHITE
        if i == 0:
            tbl.cell(i, j).fill.solid()
            tbl.cell(i, j).fill.fore_color.rgb = TJ_BLUE

# 底部说明
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE8)
box.line.color.rgb = ACCENT_GREEN
add_text_box(s, 0.8, 5.65, 12, 0.5, "▎对比公平性保证", size=16, bold=True, color=ACCENT_GREEN, font='黑体')
add_text_box(s, 0.8, 6.15, 12, 0.8,
             "所有改进方法 (Ensemble/三路融合) 均基于同一 Qwen-VL-Max baseline 进行对比，不存在模型差异。\n" +
             "本工作的提升幅度 (+3.49/+5.05/+2.94 等) 可以作为方法本身有效性的直接证据。",
             size=13, color=DARK_GRAY)

add_page_footer(s, 24, TOTAL)


# =========================================================
# Slide 25: 备用页 2 - 完整数值表
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "备用 ②  |  9 数据集完整数值表", "Backup")

tbl = s.shapes.add_table(11, 6,
                          Inches(0.3), Inches(1.0),
                          Inches(12.7), Inches(5.8)).table
hdr = ["数据集", "主指标", "Baseline", "Ensemble", "三路融合", "绝对提升"]
rows = [
    ["FIQ dress", "R@10", "15.80", "18.30", "19.29", "+3.49"],
    ["FIQ shirt", "R@10", "26.00", "27.35", "27.35", "+1.35"],
    ["FIQ toptee", "R@10", "23.09", "26.57", "27.35", "+4.26"],
    ["CIRCO", "mAP@5", "15.72", "18.06", "20.36", "+4.64"],
    ["CIRCO", "mAP@10", "16.21", "18.69", "21.26", "+5.05"],
    ["CIRR", "R@1", "22.96", "25.28", "25.90", "+2.94"],
    ["CIRR", "R@5", "53.03", "56.02", "56.28", "+3.25"],
    ["CIRR", "R@10", "65.85", "68.76", "68.72", "+2.87"],
    ["GeneCIS ch_obj", "R@1", "13.83", "—", "25.51", "+11.68"],
    ["GeneCIS fo_attr", "R@1", "18.82", "—", "27.83", "+9.01"],
]
for j, h in enumerate(hdr):
    tbl.cell(0, j).text = h
for i, r in enumerate(rows):
    for j, v in enumerate(r):
        tbl.cell(i+1, j).text = v
for i in range(11):
    for j in range(6):
        for p in tbl.cell(i, j).text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(11)
                r.font.name = '宋体'
                if i == 0:
                    r.font.bold = True; r.font.color.rgb = WHITE
                elif j == 5 and not v.startswith("—"):
                    r.font.color.rgb = ACCENT_GREEN
                    r.font.bold = True
        if i == 0:
            tbl.cell(i, j).fill.solid()
            tbl.cell(i, j).fill.fore_color.rgb = TJ_BLUE

add_text_box(s, 0.3, 6.85, 12.7, 0.3,
             "注：完整 35 项指标全部正向提升，详见论文第 4 章及 FINAL_RESULTS.md",
             size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_page_footer(s, 25, TOTAL)


# =========================================================
# Slide 26: 备用页 3 - 方法消融详情
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "备用 ③  |  各模块消融详情", "Backup")

fig_path = os.path.join(FIG_DIR, 'fig_ablation.png')
if os.path.exists(fig_path):
    s.shapes.add_picture(fig_path, Inches(0.5), Inches(1.0), width=Inches(12.0))

# 底部说明
add_text_box(s, 0.5, 4.8, 12.3, 0.5, "▎关键观察", size=18, bold=True, color=TJ_BLUE, font='黑体')
add_bullet_list(s, 0.5, 5.3, 12.3, 1.8, [
    "CIRR 上：V7 单独使用效果最好（73.5），Ensemble 反而回调（69.0）—— 但保证了全局稳定",
    "FIQ dress 上：V7 单独使用退化（13.0），Ensemble 修复（18.0）—— 证明 Ensemble 是\"救场\"设计",
    "两个场景趋势完全相反 → 没有 Ensemble 就没有 35/35 指标全部正向的结果",
], size=13)

add_page_footer(s, 26, TOTAL)


# =========================================================
# Slide 27: 备用页 4 - GeneCIS 专用 prompt 细节
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "备用 ④  |  GeneCIS 专用 Prompt 与 task-adaptive 参数", "Backup")

# 上：问题
add_text_box(s, 0.5, 1.0, 12.3, 0.5, "▎为什么 GeneCIS 需要特殊处理？", size=18, bold=True, color=TJ_RED, font='黑体')

tbl = s.shapes.add_table(4, 3,
                          Inches(0.5), Inches(1.5),
                          Inches(12.3), Inches(1.8)).table
hdr = ["特性", "FashionIQ / CIRR", "GeneCIS"]
rows = [
    ["修改文本", "具体句子 (15-30 词)", "仅 1-2 词 (如 'bus','color')"],
    ["Gallery 大小", "几千~几万", "仅 ~14 张"],
    ["检索难度", "全局匹配", "细粒度区分相似图"],
]
for j, h in enumerate(hdr): tbl.cell(0, j).text = h
for i, r in enumerate(rows):
    for j, v in enumerate(r):
        tbl.cell(i+1, j).text = v
for i in range(4):
    for j in range(3):
        for p in tbl.cell(i, j).text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(12); r.font.name = '宋体'
                if i == 0:
                    r.font.bold = True; r.font.color.rgb = WHITE
        if i == 0:
            tbl.cell(i, j).fill.solid()
            tbl.cell(i, j).fill.fore_color.rgb = TJ_RED

# 下：专用 prompt 解决方案
add_text_box(s, 0.5, 3.6, 12.3, 0.5, "▎V2 专用 prompt + task-adaptive α/β", size=18, bold=True, color=ACCENT_GREEN, font='黑体')

tbl2 = s.shapes.add_table(5, 4,
                          Inches(0.5), Inches(4.1),
                          Inches(12.3), Inches(2.5)).table
hdr2 = ["子集", "最优 β", "最优 α", "最优 R@1 (+Δ)"]
rows2 = [
    ["change_object", "0.30", "0.95", "25.51 (+11.68)"],
    ["focus_object", "1.00", "0.90", "23.62 (+7.60)"],
    ["change_attribute", "0.80", "0.80", "21.79 (+9.14)"],
    ["focus_attribute", "0.60", "0.85", "27.83 (+9.01)"],
]
for j, h in enumerate(hdr2): tbl2.cell(0, j).text = h
for i, r in enumerate(rows2):
    for j, v in enumerate(r):
        tbl2.cell(i+1, j).text = v
for i in range(5):
    for j in range(4):
        for p in tbl2.cell(i, j).text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(11); r.font.name = '宋体'
                if i == 0:
                    r.font.bold = True; r.font.color.rgb = WHITE
                elif j == 3:
                    r.font.color.rgb = ACCENT_GREEN; r.font.bold = True
        if i == 0:
            tbl2.cell(i, j).fill.solid()
            tbl2.cell(i, j).fill.fore_color.rgb = ACCENT_GREEN

add_text_box(s, 0.5, 6.7, 12.3, 0.4,
             "关键观察：focus_object β=1.0 (弃 D₂) vs change_object β=0.3 (重 D₂)——任务差异极大",
             size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

add_page_footer(s, 27, TOTAL)


# =========================================================
# Slide 28: 备用页 5 - Prompt 具体内容
# =========================================================
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "备用 ⑤  |  V7 Anti-Hallucination Prompt 具体内容", "Backup")

# 左：V7 完整原则
add_text_box(s, 0.5, 1.0, 6, 0.5, "▎V7 核心原则（完整版）", size=18, bold=True, color=TJ_BLUE, font='黑体')

box_v7 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(5.3))
box_v7.fill.solid(); box_v7.fill.fore_color.rgb = BG_LIGHT
box_v7.line.color.rgb = TJ_BLUE

prompt_lines = [
    "① Image 1: 参考图像（Reference）",
    "② Image 2: 代理图像（AI-generated, 含幻觉）",
    "",
    "【使用规则】",
    "★ 代理图仅用于「检查」修改是否正确:",
    "  - 颜色/形状是否改变",
    "  - 是否保留了应保留的",
    "  - 是否添加了不该加的",
    "",
    "★ 禁止：复制代理图细节（背景/材质/环境）",
    "★ 禁止：描述代理图中存在的内容",
    "",
    "【输出要求】",
    "- 输出简短目标描述",
    "- 长度不超过修改文本的长度",
    "- 仅保留与检索最相关的核心属性",
]
add_bullet_list(s, 0.7, 1.7, 5.6, 4.9, prompt_lines, size=11, line_spacing=1.15)

# 右：效果对比示例
add_text_box(s, 7, 1.0, 6, 0.5, "▎输出对比（change_object 示例）", size=18, bold=True, color=ACCENT_GREEN, font='黑体')

# V7 示例
add_text_box(s, 7, 1.6, 6, 0.35, "修改文本: \"donut\"", size=12, bold=True, color=DARK_GRAY)

old_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2.0), Inches(6), Inches(1.2))
old_bg.fill.solid(); old_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xE5, 0xE5)
old_bg.line.color.rgb = TJ_RED
add_text_box(s, 7.15, 2.05, 5.8, 0.3, "旧 Prompt (original)", size=10, bold=True, color=TJ_RED)
add_text_box(s, 7.15, 2.35, 5.8, 0.8,
             "\"A round donut with sprinkles on a wooden table next to a cup of coffee in a sunny café setting.\" (幻觉细节)",
             size=11, color=DARK_GRAY)

# V7
v7_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(3.4), Inches(6), Inches(1.1))
v7_bg.fill.solid(); v7_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE8)
v7_bg.line.color.rgb = ACCENT_GREEN
add_text_box(s, 7.15, 3.45, 5.8, 0.3, "V7 Anti-Hallucination", size=10, bold=True, color=ACCENT_GREEN)
add_text_box(s, 7.15, 3.75, 5.8, 0.7,
             "\"A donut on a plate on the table\" (简洁准确)",
             size=11, color=DARK_GRAY)

# V2 专用
v2_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(4.7), Inches(6), Inches(1.1))
v2_bg.fill.solid(); v2_bg.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
v2_bg.line.color.rgb = TJ_BLUE
add_text_box(s, 7.15, 4.75, 5.8, 0.3, "V2 GeneCIS 专用 (改 5-12 词要求)", size=10, bold=True, color=TJ_BLUE)
add_text_box(s, 7.15, 5.05, 5.8, 0.7,
             "\"Cinnamon-sugar donut on white plate with teapot\" (具体视觉特征)",
             size=11, color=DARK_GRAY)

# 结论
add_text_box(s, 7, 6.0, 6, 0.4, "▎结论", size=16, bold=True, color=TJ_BLUE, font='黑体')
add_text_box(s, 7, 6.4, 6, 0.5,
             "Prompt 设计本质 = 控制 MLLM 输出的信息粒度",
             size=12, color=DARK_GRAY)

add_page_footer(s, 28, TOTAL)


# =========================================================
# 保存
# =========================================================
out_file = os.path.join(OUT_DIR, '答辩PPT_杨昊明.pptx')
prs.save(out_file)
print(f"✓ PPT 已生成: {out_file}")
print(f"  总页数: {len(prs.slides)}")
